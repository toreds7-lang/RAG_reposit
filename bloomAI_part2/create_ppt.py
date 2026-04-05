# -*- coding: utf-8 -*-
"""
LangGraph 응용 3강 - 강의 슬라이드 자동 생성 (흰 배경 + 검은 글씨 템플릿)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ──────────────────────────────────────────────
# 색상 (흰 배경 / 검은 글씨 / 코드 블록만 연회색)
# ──────────────────────────────────────────────
BLACK      = RGBColor(0x00, 0x00, 0x00)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_LIGHT = RGBColor(0xF0, 0xF0, 0xF0)  # 코드 블록 배경
GRAY_MID   = RGBColor(0x99, 0x99, 0x99)  # 보조 텍스트

# ──────────────────────────────────────────────
# 프레젠테이션 초기화 (16:9)
# ──────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height


# ──────────────────────────────────────────────
# 유틸리티
# ──────────────────────────────────────────────
def new_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def tb(slide, left, top, width, height, text, size=18, bold=False,
       align=PP_ALIGN.LEFT, color=BLACK, font="맑은 고딕", anchor=MSO_ANCHOR.TOP):
    """텍스트 박스 추가 (단일 텍스트)"""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    try:
        tf.vertical_anchor = anchor
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return box


def ml(slide, left, top, width, height, lines, size=16, bold=False,
       align=PP_ALIGN.LEFT, color=BLACK, spacing=1.5, font="맑은 고딕"):
    """멀티라인 텍스트 박스"""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font
        p.alignment = align
        p.space_after = Pt(size * (spacing - 1))
    return box


def code(slide, left, top, width, height, text, size=13):
    """코드 블록 (연회색 배경 + Consolas)"""
    add_rect(slide, left, top, width, height, GRAY_LIGHT)
    box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15),
                                   width - Inches(0.4), height - Inches(0.3))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = BLACK
        p.font.name = "Consolas"
        p.space_after = Pt(2)


def title_line(slide, text, num=""):
    """슬라이드 상단 제목 + 밑줄"""
    tb(slide, Inches(0.6), Inches(0.3), Inches(11), Inches(0.7),
       text, size=28, bold=True)
    add_rect(slide, Inches(0.6), Inches(1.05), Inches(12), Pt(2), BLACK)
    if num:
        tb(slide, Inches(12), Inches(0.3), Inches(0.8), Inches(0.7),
           num, size=14, color=GRAY_MID, align=PP_ALIGN.RIGHT)


def section_slide(number, title, subtitle=""):
    """섹션 구분 슬라이드"""
    s = new_slide()
    tb(s, Inches(1), Inches(1.8), Inches(11), Inches(1),
       number, size=48, bold=True, align=PP_ALIGN.CENTER, font="Consolas")
    tb(s, Inches(1), Inches(3.2), Inches(11), Inches(1),
       title, size=36, bold=True, align=PP_ALIGN.CENTER)
    if subtitle:
        tb(s, Inches(1), Inches(4.5), Inches(11), Inches(0.7),
           subtitle, size=20, color=GRAY_MID, align=PP_ALIGN.CENTER)
    return s


# ================================================================
# 슬라이드 1: 표지
# ================================================================
s = new_slide()
tb(s, Inches(1), Inches(1.5), Inches(11), Inches(1),
   "LangGraph 응용 강의", size=48, bold=True, align=PP_ALIGN.CENTER)
tb(s, Inches(1), Inches(2.8), Inches(11), Inches(0.8),
   "3강: 워크플로우 패턴 마스터하기", size=30, bold=True, align=PP_ALIGN.CENTER)
tb(s, Inches(1), Inches(4.2), Inches(11), Inches(0.6),
   "Prompt Chaining  |  Parallelization  |  Routing", size=22, align=PP_ALIGN.CENTER)
tb(s, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
   "초보자를 위한 단계별 실습 가이드", size=18, color=GRAY_MID, align=PP_ALIGN.CENTER)


# ================================================================
# 슬라이드 2: 학습 목표
# ================================================================
s = new_slide()
title_line(s, "학습 목표", "2")
ml(s, Inches(0.8), Inches(1.5), Inches(11.5), Inches(2.5), [
    "1. Prompt Chaining  —  여러 LLM 호출을 순차적으로 연결하는 방법",
    "2. Parallelization     —  여러 작업을 동시에 병렬로 실행하는 방법",
    "3. Routing                —  조건에 따라 실행 경로를 동적으로 선택하는 방법",
], size=22, spacing=2.0)
ml(s, Inches(0.8), Inches(4.5), Inches(11.5), Inches(2.5), [
    "이 강의를 마치면 여러분은:",
    "  - LangGraph로 복잡한 AI 워크플로우를 설계할 수 있습니다",
    "  - 상황에 맞는 적절한 패턴을 선택할 수 있습니다",
    "  - State, Node, Edge 개념을 자유롭게 활용할 수 있습니다",
], size=18, spacing=1.6)


# ================================================================
# 슬라이드 3: LangGraph 핵심 개념
# ================================================================
s = new_slide()
title_line(s, "사전 지식: LangGraph 핵심 개념", "3")
tb(s, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.5),
   "LangGraph = AI 작업의 \"조립 설명서\"", size=22, bold=True)
ml(s, Inches(0.8), Inches(2.2), Inches(11.5), Inches(4.5), [
    "State (상태)",
    "  모든 노드가 공유하는 데이터 저장소",
    "  = 칠판에 적힌 메모 (모두가 읽고 쓸 수 있음)",
    "",
    "Node (노드)",
    "  데이터를 처리하는 함수",
    "  = 각 단계의 작업자 (한 가지 일을 담당)",
    "",
    "Edge (엣지)",
    "  노드 간의 연결 (실행 순서)",
    "  = 작업 순서를 정하는 화살표",
    "",
    "기본 구조:  START → Node A → Node B → END",
], size=18, spacing=1.3)


# ================================================================
# 슬라이드 4: LLM이란?
# ================================================================
s = new_slide()
title_line(s, "LLM (대규모 언어 모델) 이란?", "4")
ml(s, Inches(0.8), Inches(1.4), Inches(11.5), Inches(5.5), [
    "LLM = Large Language Model (대규모 언어 모델)",
    "",
    "쉽게 말하면: \"엄청나게 많은 글을 읽고 학습한 AI\"",
    "  - 질문에 답하기, 글 쓰기, 번역, 코드 작성 등 다양한 작업 가능",
    "  - 대표적 LLM: GPT-4, Gemini, Claude 등",
    "",
    "이 강의에서 사용하는 모델:",
], size=20, spacing=1.4)
code(s, Inches(0.8), Inches(4.5), Inches(11.5), Inches(1.5),
     'model = init_chat_model("gpt-4o-mini", temperature=1.0)\n\n'
     '# gpt-4o-mini: OpenAI의 경량 모델 (빠르고 비용 효율적)\n'
     '# temperature: 창의성 조절 (0 = 일관적, 1 = 창의적)', size=15)


# ================================================================
# 슬라이드 5: LangChain vs LangGraph
# ================================================================
s = new_slide()
title_line(s, "LangChain vs LangGraph", "5")
ml(s, Inches(0.8), Inches(1.5), Inches(5.5), Inches(2.5), [
    "[LangChain]",
    "LLM을 쉽게 사용하기 위한 도구 모음",
    "  - 모델 초기화 (init_chat_model)",
    "  - 프롬프트 관리",
    "  - 구조화 출력 (Structured Output)",
], size=18, spacing=1.4)
ml(s, Inches(7.0), Inches(1.5), Inches(5.5), Inches(2.5), [
    "[LangGraph]",
    "LLM 작업을 그래프로 설계하는 엔진",
    "  - StateGraph로 워크플로우 구성",
    "  - 순차/병렬/조건부 실행 지원",
    "  - 복잡한 AI 파이프라인 관리",
], size=18, spacing=1.4)
ml(s, Inches(0.8), Inches(5.0), Inches(11.5), Inches(2.0), [
    "쉬운 비유:",
    "  LangChain = 요리 재료와 도구 (칼, 냄비, 양념 등)",
    "  LangGraph = 요리 레시피 (\"먼저 재료를 썰고 → 볶고 → 양념하기\")",
    "  둘을 합치면 = 완성된 요리 (AI 워크플로우)!",
], size=18, spacing=1.5)


# ================================================================
# 슬라이드 6: StateGraph 핵심 구조
# ================================================================
s = new_slide()
title_line(s, "StateGraph 핵심 구조 — 4단계로 만들기", "6")
ml(s, Inches(0.8), Inches(1.4), Inches(11.5), Inches(1.5), [
    "1. State 정의  →  TypedDict로 공유 데이터 구조 선언",
    "2. Node 정의   →  State를 처리하는 함수 작성",
    "3. Edge 연결   →  add_edge / add_conditional_edges",
    "4. Compile       →  .compile()로 실행 가능하게 변환",
], size=20, spacing=1.6)
code(s, Inches(0.8), Inches(3.8), Inches(11.5), Inches(3.0),
     '# 이 4단계가 모든 패턴의 기본 골격입니다!\n'
     'app = StateGraph(MyState)        # State 스키마 지정\n'
     'app.add_node(\'name\', function)   # 노드 등록\n'
     'app.add_edge(START, \'name\')      # 엣지 연결\n'
     'chain = app.compile()            # 컴파일\n'
     'result = chain.invoke(inputs)    # 실행!', size=16)


# ================================================================
# 슬라이드 7: TypedDict & Pydantic
# ================================================================
s = new_slide()
title_line(s, "TypedDict & Pydantic — 데이터 타입 정의", "7")
tb(s, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.4),
   "TypedDict — State 정의에 사용", size=20, bold=True)
ml(s, Inches(0.8), Inches(1.8), Inches(5.5), Inches(0.8), [
    "- 딕셔너리의 각 키에 타입을 지정",
    "- 어떤 데이터가 오가는지 명확하게 선언",
], size=16, spacing=1.4)
code(s, Inches(0.8), Inches(2.8), Inches(5.5), Inches(1.6),
     'class JokeState(TypedDict):\n'
     '    topic: str        # 문자열\n'
     '    draft_joke: str   # 문자열\n'
     '    final_joke: str   # 문자열', size=14)

tb(s, Inches(7.0), Inches(1.3), Inches(5.5), Inches(0.4),
   "Pydantic — 구조화 출력에 사용", size=20, bold=True)
ml(s, Inches(7.0), Inches(1.8), Inches(5.5), Inches(0.8), [
    "- LLM 출력을 정해진 형식으로 강제",
    "- 데이터 검증 기능 내장",
], size=16, spacing=1.4)
code(s, Inches(7.0), Inches(2.8), Inches(5.5), Inches(1.6),
     "class RouteDecision(BaseModel):\n"
     "    category: Literal[\n"
     "        'billing', 'technical',\n"
     "        'shipping', 'general'\n"
     "    ]", size=14)

ml(s, Inches(0.8), Inches(5.0), Inches(11.5), Inches(1.5), [
    "요약:",
    "  TypedDict = \"이 칠판에는 이런 정보가 적힌다\" 선언",
    "  Pydantic = \"AI야, 반드시 이 형식으로만 대답해\" 강제",
], size=18, spacing=1.5)


# ================================================================
# 슬라이드 8: Prompt Chaining 섹션
# ================================================================
section_slide("PATTERN 1", "Prompt Chaining (프롬프트 체이닝)",
              "여러 LLM 호출을 순차적으로 연결하는 파이프라인")


# ================================================================
# 슬라이드 9: Prompt Chaining 개념
# ================================================================
s = new_slide()
title_line(s, "Prompt Chaining 이란?", "9")
ml(s, Inches(0.8), Inches(1.4), Inches(11.5), Inches(2.0), [
    "이전 단계의 출력이 다음 단계의 입력이 되는 순차적 파이프라인",
    "",
    "  - 복잡한 작업을 작은 단계로 분리 → 품질 향상",
    "  - 각 단계마다 다른 지시사항(프롬프트)을 줄 수 있음",
], size=20, spacing=1.4)
ml(s, Inches(0.8), Inches(3.5), Inches(11.5), Inches(1.0), [
    "흐름:  START → [1단계: 초안 생성] → [2단계: 수정/개선] → [3단계: 최종 포장] → END",
], size=20, bold=True, spacing=1.5)
ml(s, Inches(0.8), Inches(4.8), Inches(11.5), Inches(1.5), [
    "비유: 제빵 과정 → 반죽 만들기 → 모양 잡기 → 굽기 → 장식하기",
    "",
    "이번 예제: 주제 → 농담 초안 생성 → 아재개그로 개선 → 이모지 추가(SNS용)",
], size=18, spacing=1.5)


# ================================================================
# 슬라이드 10: State 정의
# ================================================================
s = new_slide()
title_line(s, "핵심 코드 1: State 정의 (JokeState)", "10")
ml(s, Inches(0.8), Inches(1.4), Inches(11.5), Inches(0.8), [
    "파이프라인의 각 단계에서 사용할 데이터를 미리 선언합니다.",
    "각 노드는 이 State를 읽고, 자기 결과를 저장합니다.",
], size=18, spacing=1.5)
code(s, Inches(0.8), Inches(2.6), Inches(11.5), Inches(2.3),
     'class JokeState(TypedDict):\n'
     '    topic: str          # 입력: 농담의 주제 (예: "고양이")\n'
     '    draft_joke: str     # 1단계 결과: AI가 생성한 농담 초안\n'
     '    improved_joke: str  # 2단계 결과: 아재개그 스타일로 개선된 버전\n'
     '    final_joke: str     # 3단계 결과: 이모지가 추가된 최종 완성본', size=16)
ml(s, Inches(0.8), Inches(5.3), Inches(11.5), Inches(1.5), [
    "핵심 포인트:",
    "  - State는 \"칠판\"이다 → 모든 노드가 여기에 읽고 쓴다",
    "  - 각 필드 = 칠판의 한 칸 → 어떤 데이터가 올지 미리 약속",
], size=17, spacing=1.5)


# ================================================================
# 슬라이드 11: Node 함수
# ================================================================
s = new_slide()
title_line(s, "핵심 코드 2: Node 함수 (처리 단계)", "11")
tb(s, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.4),
   "각 노드는 State에서 데이터를 읽고 → LLM 호출 → 결과를 State에 저장", size=18)

tb(s, Inches(0.8), Inches(1.8), Inches(5.5), Inches(0.3),
   "1단계: 초안 생성", size=16, bold=True)
code(s, Inches(0.8), Inches(2.2), Inches(5.7), Inches(2.3),
     'def generate_joke(state: JokeState):\n'
     '    topic = state["topic"]       # State에서 읽기\n'
     '    msg = model.invoke(\n'
     '        f"\\\'{ topic }\\\'에 대한 농담을 만들어줘."\n'
     '    )\n'
     '    return {"draft_joke": msg.content}  # State에 쓰기', size=13)

tb(s, Inches(6.8), Inches(1.8), Inches(5.5), Inches(0.3),
   "2단계: 아재개그로 개선", size=16, bold=True)
code(s, Inches(6.8), Inches(2.2), Inches(5.7), Inches(2.3),
     'def critique_and_improve(state: JokeState):\n'
     '    original = state["draft_joke"]  # 이전 결과 읽기\n'
     '    prompt = f"다음 농담을 아재개그로 개선해줘.\n'
     '    원문: {original}"\n'
     '    msg = model.invoke(prompt)\n'
     '    return {"improved_joke": msg.content}', size=13)

ml(s, Inches(0.8), Inches(4.9), Inches(11.5), Inches(2.0), [
    "핵심 패턴:",
    "  state[\"키\"]로 읽기  →  model.invoke()로 LLM 호출  →  {\"키\": 결과}로 반환",
    "",
    "  반환하지 않은 키는 그대로 유지됩니다! (topic은 계속 살아있음)",
    "  3단계(polish_joke)도 동일한 패턴: improved_joke 읽기 → 이모지 추가 → final_joke 반환",
], size=16, spacing=1.3)


# ================================================================
# 슬라이드 12: 그래프 생성
# ================================================================
s = new_slide()
title_line(s, "핵심 코드 3: 그래프 생성 & 컴파일", "12")
code(s, Inches(0.8), Inches(1.4), Inches(11.5), Inches(4.5),
     '# 1. StateGraph 생성 - State 스키마 지정\n'
     'app = StateGraph(JokeState)\n'
     '\n'
     '# 2. 노드 등록 - (이름, 함수)\n'
     'app.add_node(\'generate_joke\', generate_joke)\n'
     'app.add_node(\'critique_and_improve\', critique_and_improve)\n'
     'app.add_node(\'polish_joke\', polish_joke)\n'
     '\n'
     '# 3. 엣지 연결 - 실행 순서 지정\n'
     'app.add_edge(START, \'generate_joke\')                  # 시작 → 1단계\n'
     'app.add_edge(\'generate_joke\', \'critique_and_improve\') # 1단계 → 2단계\n'
     'app.add_edge(\'critique_and_improve\', \'polish_joke\')   # 2단계 → 3단계\n'
     'app.add_edge(\'polish_joke\', END)                      # 3단계 → 종료\n'
     '\n'
     '# 4. 컴파일 - 실행 가능하게 변환\n'
     'chain = app.compile()', size=14)
ml(s, Inches(0.8), Inches(6.1), Inches(11.5), Inches(1.0), [
    "START와 END는 LangGraph의 특별한 예약 노드입니다.",
    "add_edge()는 \"A 다음에 B를 실행하라\"는 뜻!",
], size=17, spacing=1.4)


# ================================================================
# 슬라이드 13: 실행 & 정리
# ================================================================
s = new_slide()
title_line(s, "실행 & 정리: Prompt Chaining", "13")
code(s, Inches(0.8), Inches(1.4), Inches(11.5), Inches(2.5),
     '# 실행: 초기 State에 topic만 넣어주면 됩니다!\n'
     'result = chain.invoke({"topic": "고양이"})\n'
     '\n'
     '# 결과 확인\n'
     'print(result[\'draft_joke\'])     # 1단계: 초안\n'
     'print(result[\'improved_joke\'])  # 2단계: 아재개그\n'
     'print(result[\'final_joke\'])     # 3단계: SNS용 완성본', size=15)
ml(s, Inches(0.8), Inches(4.3), Inches(11.5), Inches(2.5), [
    "Prompt Chaining 핵심 정리:",
    "",
    "  구조: A → B → C (순차 실행, Linear Chain)",
    "  장점: 각 단계별 품질 제어 가능, 디버깅 용이",
    "  사용처: 초안→검토→최종본, 번역→교정→포맷팅 등",
    "  핵심 API: add_edge(A, B) — 순차 연결",
], size=18, spacing=1.4)


# ================================================================
# 슬라이드 14: Parallelization 섹션
# ================================================================
section_slide("PATTERN 2", "Parallelization (병렬 처리)",
              "여러 작업을 동시에 병렬로 실행하는 패턴")


# ================================================================
# 슬라이드 15: Parallelization 개념
# ================================================================
s = new_slide()
title_line(s, "Parallelization 이란?", "15")
ml(s, Inches(0.8), Inches(1.4), Inches(11.5), Inches(1.5), [
    "서로 독립적인 작업을 동시에 실행하여 시간을 절약하는 패턴",
    "",
    "핵심 용어:  Fan-out (분기: 1→여럿)  /  Fan-in (합류: 여럿→1)",
], size=20, spacing=1.4)
ml(s, Inches(0.8), Inches(3.2), Inches(11.5), Inches(3.0), [
    "               ┌→ [Worker A: 시 작성]   ─────┐",
    "START → Router ┼→ [Worker B: 소설 작성] ──┼→ [Aggregator: 취합] → END",
    "               └→ [Worker C: 농담 작성] ──┘",
    "",
    "         Fan-out (분기)                     Fan-in (합류)",
], size=18, spacing=1.3, font="Consolas")
ml(s, Inches(0.8), Inches(5.8), Inches(11.5), Inches(1.0), [
    "비유: 식당에서 요리사 3명이 동시에 각자 요리 → 한 접시에 담기",
    "이번 예제: 하나의 주제로 시, 소설, 농담을 동시에 생성하고 하나의 리포트로 합침",
], size=18, spacing=1.5)


# ================================================================
# 슬라이드 16: WriterState
# ================================================================
s = new_slide()
title_line(s, "핵심 코드 1: State 정의 (WriterState)", "16")
tb(s, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.4),
   "각 Worker가 독립적으로 결과를 저장할 필드를 선언합니다.", size=18)
code(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(2.8),
     'class WriterState(TypedDict):\n'
     '    topic: str          # 입력: 글쓰기 주제\n'
     '    poem: str           # Worker A 결과: 시(Poem)\n'
     '    story: str          # Worker B 결과: 소설(Story)\n'
     '    joke: str           # Worker C 결과: 농담(Joke)\n'
     '    final_report: str   # Aggregator 결과: 최종 편집본', size=16)
ml(s, Inches(0.8), Inches(5.2), Inches(11.5), Inches(1.8), [
    "핵심 포인트:",
    "  - 각 Worker는 자기 필드에만 쓴다 (poem, story, joke)",
    "  - 서로 다른 필드에 쓰기 때문에 충돌 없이 동시 실행 가능!",
    "  - Aggregator가 모든 결과를 읽어서 final_report에 합침",
], size=17, spacing=1.4)


# ================================================================
# 슬라이드 17: Worker & Aggregator
# ================================================================
s = new_slide()
title_line(s, "핵심 코드 2: Worker & Aggregator 노드", "17")

tb(s, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.3),
   "Worker 노드 (3개 모두 같은 패턴)", size=16, bold=True)
code(s, Inches(0.8), Inches(1.7), Inches(5.7), Inches(2.5),
     '# Worker A: 시 작성 (B, C도 동일 패턴!)\n'
     'def write_poem(state: WriterState):\n'
     '    topic = state["topic"]\n'
     '    msg = model.invoke(\n'
     '        f"\\\'{ topic }\\\'에 대한 시를 짧게 써줘."\n'
     '    )\n'
     '    return {"poem": msg.content}', size=13)

tb(s, Inches(6.8), Inches(1.3), Inches(5.5), Inches(0.3),
   "Aggregator 노드 (Fan-in)", size=16, bold=True)
code(s, Inches(6.8), Inches(1.7), Inches(5.7), Inches(2.5),
     '# Aggregator: 모든 결과 합치기\n'
     'def aggregator(state: WriterState):\n'
     '    final = f"""\n'
     "    1. 시: {state['poem']}\n"
     "    2. 소설: {state['story']}\n"
     "    3. 농담: {state['joke']}\n"
     '    """\n'
     '    return {"final_report": final}', size=13)

ml(s, Inches(0.8), Inches(4.5), Inches(11.5), Inches(2.5), [
    "핵심 패턴:",
    "  - Worker: topic 읽기 → LLM 호출 → 자기 필드에 저장",
    "  - Aggregator: 모든 Worker의 결과를 State에서 읽어서 합침",
    "  - LangGraph가 자동으로 모든 Worker 완료를 기다린 후 Aggregator 실행!",
], size=17, spacing=1.4)


# ================================================================
# 슬라이드 18: Fan-out / Fan-in 그래프
# ================================================================
s = new_slide()
title_line(s, "핵심 코드 3: 그래프 생성 (Fan-out / Fan-in)", "18")
code(s, Inches(0.8), Inches(1.4), Inches(7.0), Inches(5.3),
     'workflow = StateGraph(WriterState)\n'
     '\n'
     '# 노드 등록\n'
     "workflow.add_node('write_poem', write_poem)\n"
     "workflow.add_node('write_story', write_story)\n"
     "workflow.add_node('write_joke', write_joke)\n"
     "workflow.add_node('aggregator', aggregator)\n"
     '\n'
     '# Fan-out: START에서 3개 Worker로 동시 분기!\n'
     "workflow.add_edge(START, 'write_poem')\n"
     "workflow.add_edge(START, 'write_story')\n"
     "workflow.add_edge(START, 'write_joke')\n"
     '\n'
     '# Fan-in: 3개 Worker에서 Aggregator로 합류!\n'
     "workflow.add_edge('write_poem', 'aggregator')\n"
     "workflow.add_edge('write_story', 'aggregator')\n"
     "workflow.add_edge('write_joke', 'aggregator')\n"
     '\n'
     "workflow.add_edge('aggregator', END)\n"
     'app = workflow.compile()', size=13)

ml(s, Inches(8.2), Inches(1.4), Inches(4.3), Inches(2.5), [
    "Fan-out (분기):",
    "START에서 여러 노드로",
    "동시에 엣지를 연결하면",
    "LangGraph가 자동으로",
    "병렬 실행합니다!",
], size=16, spacing=1.3)
ml(s, Inches(8.2), Inches(4.0), Inches(4.3), Inches(2.5), [
    "Fan-in (합류):",
    "여러 노드에서 하나로",
    "엣지를 연결하면",
    "모든 노드 완료 후",
    "자동으로 다음 실행!",
], size=16, spacing=1.3)


# ================================================================
# 슬라이드 19: 실행 & 정리
# ================================================================
s = new_slide()
title_line(s, "실행 & 정리: Parallelization", "19")
code(s, Inches(0.8), Inches(1.4), Inches(11.5), Inches(1.6),
     '# 실행\n'
     'result = app.invoke({"topic": "직장인의 월요일"})\n'
     'print(result["final_report"])', size=15)
ml(s, Inches(0.8), Inches(3.3), Inches(11.5), Inches(3.5), [
    "순차 실행 vs 병렬 실행 비교:",
    "",
    "  순차: 시 작성(3초) → 소설 작성(3초) → 농담 작성(3초) = 총 9초",
    "  병렬: 시 + 소설 + 농담 동시 작성 = 총 3초! (3배 빠름)",
    "",
    "핵심: START에서 여러 노드로 엣지를 연결하면 자동 병렬 실행!",
], size=20, spacing=1.5)


# ================================================================
# 슬라이드 20: Routing 섹션
# ================================================================
section_slide("PATTERN 3", "Routing (라우팅)",
              "조건에 따라 실행 경로를 동적으로 선택하는 패턴")


# ================================================================
# 슬라이드 21: Routing 개념
# ================================================================
s = new_slide()
title_line(s, "Routing 이란?", "21")
ml(s, Inches(0.8), Inches(1.4), Inches(11.5), Inches(1.0), [
    "입력 내용을 분석하여 적절한 처리 경로를 자동으로 선택하는 패턴",
    "핵심: 조건부 엣지 (Conditional Edge) — 함수 반환값에 따라 다음 노드 결정",
], size=20, spacing=1.5)
ml(s, Inches(0.8), Inches(3.0), Inches(11.5), Inches(2.5), [
    "              ┌─ 결제 문의 → [Billing Expert]",
    "              ├─ 기술 문의 → [Tech Expert]",
    "START → Router ┤",
    "              ├─ 배송 문의 → [Shipping Expert]",
    "              └─ 기타 문의 → [General Expert]",
], size=18, spacing=1.2, font="Consolas")
ml(s, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.5), [
    "비유: 병원 접수대 → 증상에 따라 내과/외과/안과로 자동 안내",
    "이번 예제: 고객 문의 내용을 LLM이 분석하여 적합한 담당 부서로 자동 연결",
], size=18, spacing=1.5)


# ================================================================
# 슬라이드 22: Structured Output
# ================================================================
s = new_slide()
title_line(s, "핵심 코드 1: Structured Output (구조화 출력)", "22")
ml(s, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.8), [
    "LLM이 자유 텍스트 대신 정해진 형식으로만 응답하도록 강제합니다.",
    "이를 통해 라우터가 정확한 카테고리 값을 받을 수 있습니다.",
], size=18, spacing=1.5)
code(s, Inches(0.8), Inches(2.5), Inches(11.5), Inches(3.0),
     'from typing import Literal\n'
     '\n'
     '# 출력 스키마 정의 - LLM은 이 형식으로만 응답 가능!\n'
     'class RouteDecision(BaseModel):\n'
     "    category: Literal['billing', 'technical', 'shipping', 'general'] = Field(\n"
     "        description='고객 문의를 분석하여 적절한 부서를 선택하세요.'\n"
     '    )\n'
     '\n'
     '# 구조화 출력을 사용하는 라우터 LLM 생성\n'
     'router_llm = model.with_structured_output(RouteDecision)', size=15)
ml(s, Inches(0.8), Inches(5.8), Inches(11.5), Inches(1.2), [
    "Literal['billing', 'technical', ...] = 이 4개 값 중 하나만 선택 가능",
    "with_structured_output() = \"AI야, 반드시 RouteDecision 형식으로만 답해!\"",
], size=17, spacing=1.5)


# ================================================================
# 슬라이드 23: Router & Expert 노드
# ================================================================
s = new_slide()
title_line(s, "핵심 코드 2: Router 노드 & Expert 노드", "23")

tb(s, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.3),
   "Router 노드", size=16, bold=True)
code(s, Inches(0.8), Inches(1.7), Inches(5.7), Inches(2.3),
     '# Router: 문의 분석 → 카테고리 결정\n'
     'def router_node(state: SupportState):\n'
     '    query = state["query"]\n'
     '    # 구조화 출력 LLM 호출\n'
     '    decision = router_llm.invoke(query)\n'
     '    # decision.category = \'billing\' 등\n'
     '    return {"category": decision.category}', size=13)

tb(s, Inches(6.8), Inches(1.3), Inches(5.5), Inches(0.3),
   "Expert 노드 (4개)", size=16, bold=True)
code(s, Inches(6.8), Inches(1.7), Inches(5.7), Inches(2.3),
     '# Expert 노드 (4개 중 하나만 실행됨!)\n'
     'def billing_expert(state: SupportState):\n'
     '    prompt = f"결제 전문가로 답변: {state[\'query\']}"\n'
     '    msg = model.invoke(prompt)\n'
     '    return {"response": msg.content}\n'
     '\n'
     '# technical, shipping, general도 동일', size=13)

ml(s, Inches(0.8), Inches(4.3), Inches(11.5), Inches(2.8), [
    "실행 흐름:",
    "",
    "  1. 고객 문의 입력 → Router가 LLM으로 카테고리 분류",
    "  2. 분류 결과(category)를 State에 저장",
    "  3. 조건부 엣지가 category 값을 읽어서 해당 Expert만 실행!",
    "  4. Expert가 전문 답변 생성 → response에 저장 → 종료",
], size=17, spacing=1.3)


# ================================================================
# 슬라이드 24: 조건부 엣지
# ================================================================
s = new_slide()
title_line(s, "핵심 코드 3: 조건부 엣지 (Conditional Edge)", "24")
code(s, Inches(0.8), Inches(1.4), Inches(11.5), Inches(4.0),
     '# 라우팅 함수: State의 category 값에 따라 다음 노드 결정\n'
     'def route_to_expert(state: SupportState):\n'
     "    category = state['category']\n"
     "    if category == 'billing':    return 'billing_expert'\n"
     "    elif category == 'technical': return 'technical_expert'\n"
     "    elif category == 'shipping':  return 'shipping_expert'\n"
     "    else:                         return 'general_expert'\n"
     '\n'
     '# 그래프에 조건부 엣지 등록\n'
     'workflow.add_conditional_edges(\n'
     "    'router_node',       # 이 노드 실행 후 분기\n"
     '    route_to_expert,     # 분기를 결정하는 함수\n'
     "    ['billing_expert', 'technical_expert',     # 가능한 다음 노드\n"
     "     'shipping_expert', 'general_expert']\n"
     ')', size=14)
ml(s, Inches(0.8), Inches(5.7), Inches(11.5), Inches(1.3), [
    "add_edge() vs add_conditional_edges():",
    "  add_edge(A, B) = \"A 다음은 항상 B\" (고정)",
    "  add_conditional_edges(A, func, [...]) = \"A 다음은 func()의 결과에 따라 결정\" (동적)",
], size=17, spacing=1.4)


# ================================================================
# 슬라이드 25: 실행 & 정리
# ================================================================
s = new_slide()
title_line(s, "실행 & 정리: Routing", "25")
code(s, Inches(0.8), Inches(1.4), Inches(11.5), Inches(3.2),
     '# 테스트 1: 결제 문의 → billing 부서로 자동 라우팅\n'
     'result = app.invoke({"query": "지난달 요금이 두 번 빠져나갔어요."})\n'
     "print(result['category'])   # → 'billing'\n"
     '\n'
     '# 테스트 2: 기술 문의 → technical 부서로 자동 라우팅\n'
     'result = app.invoke({"query": "API 연결할 때 404 에러가 자꾸 떠요."})\n'
     "print(result['category'])   # → 'technical'\n"
     '\n'
     '# 테스트 3: 배송 문의 → shipping 부서로 자동 라우팅\n'
     'result = app.invoke({"query": "주문한 노트북 언제 도착하나요?"})\n'
     "print(result['category'])   # → 'shipping'", size=14)
ml(s, Inches(0.8), Inches(5.0), Inches(11.5), Inches(2.0), [
    "Routing 핵심 정리:",
    "  구조: Router → 조건에 따라 A 또는 B 또는 C (동적 분기)",
    "  핵심 API: add_conditional_edges() + Structured Output",
    "  사용처: 고객 지원 분류, 질문 유형별 처리, 콘텐츠 라우팅 등",
], size=18, spacing=1.4)


# ================================================================
# 슬라이드 26: 3가지 패턴 비교표
# ================================================================
s = new_slide()
title_line(s, "3가지 패턴 비교표", "26")
ml(s, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5), [
    "┌──────────────┬────────────────────┬─────────────────────┬────────────────────┐",
    "│              │ Prompt Chaining    │ Parallelization     │ Routing            │",
    "├──────────────┼────────────────────┼─────────────────────┼────────────────────┤",
    "│ 구조         │ A → B → C          │ A, B, C 동시 → D   │ 조건 → A or B or C │",
    "│              │ (순차 실행)        │ (Fan-out / Fan-in)  │ (동적 분기)        │",
    "├──────────────┼────────────────────┼─────────────────────┼────────────────────┤",
    "│ 장점         │ 단계적 품질 향상   │ 처리 시간 단축      │ 적재적소 처리      │",
    "│              │ 디버깅 용이        │ 독립 작업 동시 실행 │ 동적 경로 선택     │",
    "├──────────────┼────────────────────┼─────────────────────┼────────────────────┤",
    "│ 사용 사례    │ 초안→검토→최종본   │ 다각도 분석         │ 고객 지원 분류     │",
    "│              │ 번역→교정→포맷팅   │ 다중 콘텐츠 생성    │ 질문 유형별 처리   │",
    "├──────────────┼────────────────────┼─────────────────────┼────────────────────┤",
    "│ 핵심 API     │ add_edge(A, B)     │ add_edge(START, A)  │ add_conditional_   │",
    "│              │                    │ add_edge(START, B)  │ edges(A, func, []) │",
    "└──────────────┴────────────────────┴─────────────────────┴────────────────────┘",
], size=14, spacing=1.0, font="Consolas")


# ================================================================
# 슬라이드 27: 핵심 개념 총정리
# ================================================================
s = new_slide()
title_line(s, "핵심 개념 총정리", "27")
ml(s, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5), [
    "State",
    "  그래프 전체에서 공유되는 데이터 저장소 (TypedDict로 정의)",
    "",
    "Node",
    "  상태를 처리하는 함수 (입력: State, 출력: 업데이트할 딕셔너리)",
    "",
    "Edge",
    "  노드 간 연결 — add_edge(일반) vs add_conditional_edges(조건부)",
    "",
    "Fan-out",
    "  START → 여러 노드 : 병렬 실행 시작",
    "",
    "Fan-in",
    "  여러 노드 → 하나의 노드 : 결과 취합 (자동 대기)",
    "",
    "Structured Output",
    "  LLM 출력을 Pydantic 모델로 강제 (with_structured_output)",
], size=18, spacing=1.1)


# ================================================================
# 슬라이드 28: Q&A
# ================================================================
s = new_slide()
tb(s, Inches(1), Inches(2.0), Inches(11), Inches(1.2),
   "Q & A", size=60, bold=True, align=PP_ALIGN.CENTER, font="Consolas")
tb(s, Inches(1), Inches(3.5), Inches(11), Inches(0.8),
   "질문이 있으신가요?", size=28, align=PP_ALIGN.CENTER)
ml(s, Inches(1), Inches(5.0), Inches(11), Inches(1.5), [
    "오늘 배운 3가지 패턴을 직접 실습해보세요!",
    "주제를 바꿔가며 코드를 실행해보면 이해가 더 깊어집니다.",
], size=20, spacing=1.5, align=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────
# 저장
# ──────────────────────────────────────────────
output_path = r"c:\DL_Practice\[3강]_랭그래프_응용_강의슬라이드.pptx"
prs.save(output_path)
print(f"PPT 생성 완료: {output_path}")
print(f"총 슬라이드 수: {len(prs.slides)}")
