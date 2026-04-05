# -*- coding: utf-8 -*-
"""
LangGraph 응용 3강 (Part 2) - 강의 슬라이드 자동 생성 (흰 배경 + 검은 글씨 템플릿)
Evaluator-Optimizer | Orchestrator-Worker | Agents
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ──────────────────────────────────────────────
# 색상 (흰 배경 / 검은 글씨 / 코드 블록만 연회색)
# ──────────────────────────────────────────────
BLACK      = RGBColor(0x00, 0x00, 0x00)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_LIGHT = RGBColor(0xF0, 0xF0, 0xF0)
GRAY_MID   = RGBColor(0x99, 0x99, 0x99)

# ──────────────────────────────────────────────
# 프레젠테이션 초기화 (16:9)
# ──────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height


# ──────────────────────────────────────────────
# 유틸리티
# ──────────────────────────────────────────────
def new_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def tb(slide, left, top, width, height, text, size=18, bold=False,
       align=PP_ALIGN.LEFT, color=BLACK, font="맑은 고딕", anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    try:
        tf.vertical_anchor = anchor
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return box


def ml(slide, left, top, width, height, lines, size=16, bold=False,
       align=PP_ALIGN.LEFT, color=BLACK, spacing=1.5, font="맑은 고딕"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font
        p.alignment = align
        p.space_after = Pt(size * (spacing - 1))
    return box


def code(slide, left, top, width, height, text, size=13):
    add_rect(slide, left, top, width, height, GRAY_LIGHT)
    box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15),
                                   width - Inches(0.4), height - Inches(0.3))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = BLACK
        p.font.name = "Consolas"
        p.space_after = Pt(2)


def title_line(slide, text, num=""):
    tb(slide, Inches(0.6), Inches(0.3), Inches(11), Inches(0.7),
       text, size=28, bold=True)
    add_rect(slide, Inches(0.6), Inches(1.05), Inches(12), Pt(2), BLACK)
    if num:
        tb(slide, Inches(12), Inches(0.3), Inches(0.8), Inches(0.7),
           num, size=14, color=GRAY_MID, align=PP_ALIGN.RIGHT)


def section_slide(number, title, subtitle=""):
    s = new_slide()
    tb(s, Inches(1), Inches(1.8), Inches(11), Inches(1),
       number, size=48, bold=True, align=PP_ALIGN.CENTER, font="Consolas")
    tb(s, Inches(1), Inches(3.2), Inches(11), Inches(1),
       title, size=36, bold=True, align=PP_ALIGN.CENTER)
    if subtitle:
        tb(s, Inches(1), Inches(4.5), Inches(11), Inches(0.7),
           subtitle, size=20, color=GRAY_MID, align=PP_ALIGN.CENTER)
    return s


# ================================================================
# 슬라이드 1: 표지
# ================================================================
s = new_slide()
tb(s, Inches(1), Inches(1.5), Inches(11), Inches(1),
   "LangGraph 응용 강의", size=48, bold=True, align=PP_ALIGN.CENTER)
tb(s, Inches(1), Inches(2.8), Inches(11), Inches(0.8),
   "3강 (Part 2): 고급 워크플로우 패턴", size=30, bold=True, align=PP_ALIGN.CENTER)
tb(s, Inches(1), Inches(4.2), Inches(11), Inches(0.6),
   "Evaluator-Optimizer  |  Orchestrator-Worker  |  Agents", size=22, align=PP_ALIGN.CENTER)
tb(s, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
   "초보자를 위한 단계별 실습 가이드", size=18, color=GRAY_MID, align=PP_ALIGN.CENTER)


# ================================================================
# 슬라이드 2: 학습 목표
# ================================================================
s = new_slide()
title_line(s, "학습 목표", "2")
ml(s, Inches(0.8), Inches(1.5), Inches(11.5), Inches(2.5), [
    "1. Evaluator-Optimizer  —  생성과 평가를 반복하여 결과물을 개선하는 루프 패턴",
    "2. Orchestrator-Worker  —  큰 작업을 분할하여 병렬로 처리하는 Map-Reduce 패턴",
    "3. Agents (에이전트)       —  LLM이 스스로 도구를 선택하고 실행하는 자율 패턴",
], size=22, spacing=2.0)
ml(s, Inches(0.8), Inches(4.5), Inches(11.5), Inches(2.5), [
    "이 강의를 마치면 여러분은:",
    "  - 조건부 엣지로 루프(반복) 워크플로우를 설계할 수 있습니다",
    "  - Send API로 동적 병렬 처리를 구현할 수 있습니다",
    "  - LLM 에이전트에 도구(Tool)를 연결할 수 있습니다",
], size=18, spacing=1.6)


# ================================================================
# 슬라이드 3: LangGraph 핵심 개념 복습
# ================================================================
s = new_slide()
title_line(s, "복습: LangGraph 핵심 개념", "3")
ml(s, Inches(0.8), Inches(1.4), Inches(11.5), Inches(5.5), [
    "State (상태)",
    "  그래프 전체에서 공유되는 데이터 저장소 (TypedDict로 정의)",
    "  = 모든 노드가 읽고 쓰는 칠판",
    "",
    "Node (노드)",
    "  State를 입력받아 처리하고, 변경된 부분만 딕셔너리로 반환하는 함수",
    "",
    "Edge (엣지)",
    "  add_edge(A, B) : A 다음에 항상 B 실행 (고정 경로)",
    "  add_conditional_edges(A, func, [...]) : func() 결과에 따라 동적 분기",
    "",
    "기본 흐름:  START -> Node A -> Node B -> END",
    "루프 흐름:  START -> A -> B -> (조건) -> A 또는 END",
], size=18, spacing=1.2)


# ================================================================
# 슬라이드 4: Structured Output & Pydantic
# ================================================================
s = new_slide()
title_line(s, "복습: Structured Output & Pydantic", "4")
ml(s, Inches(0.8), Inches(1.4), Inches(11.5), Inches(1.5), [
    "LLM은 기본적으로 자유로운 텍스트를 생성합니다.",
    "하지만 프로그램에서 사용하려면 정해진 형식이 필요합니다!",
], size=20, spacing=1.5)

tb(s, Inches(0.8), Inches(2.8), Inches(5.5), Inches(0.4),
   "Pydantic BaseModel로 형식 정의", size=18, bold=True)
code(s, Inches(0.8), Inches(3.3), Inches(5.5), Inches(2.0),
     "class EvaluationResult(BaseModel):\n"
     "    status: Literal['pass', 'fail']\n"
     "    feedback: str\n"
     "\n"
     "# LLM은 이 형식으로만 응답 가능!", size=14)

tb(s, Inches(7.0), Inches(2.8), Inches(5.5), Inches(0.4),
   "with_structured_output으로 적용", size=18, bold=True)
code(s, Inches(7.0), Inches(3.3), Inches(5.5), Inches(2.0),
     "evaluator_llm = model.with_structured_output(\n"
     "    EvaluationResult\n"
     ")\n"
     "result = evaluator_llm.invoke(prompt)\n"
     "# result.status = 'pass' or 'fail'", size=14)

ml(s, Inches(0.8), Inches(5.7), Inches(11.5), Inches(1.3), [
    "핵심: Pydantic = \"AI야, 반드시 이 형식으로만 답해!\"",
    "  Literal['A', 'B'] = 지정된 값 중 하나만 선택 가능",
], size=18, spacing=1.5)


# ================================================================
# 슬라이드 5: 조건부 엣지와 루프
# ================================================================
s = new_slide()
title_line(s, "복습: 조건부 엣지로 루프 만들기", "5")
ml(s, Inches(0.8), Inches(1.4), Inches(11.5), Inches(1.5), [
    "조건부 엣지 = \"상황에 따라 다른 경로로 가라\"",
    "이전 노드로 다시 돌아가면 루프(반복)가 됩니다!",
], size=20, spacing=1.5)
code(s, Inches(0.8), Inches(3.0), Inches(11.5), Inches(2.5),
     '# 라우팅 함수: 다음에 어디로 갈지 결정\n'
     'def route_function(state):\n'
     '    if state["status"] == "pass":   return END            # 합격 -> 종료\n'
     '    if state["count"] >= 3:         return END            # 3번 시도 -> 강제 종료\n'
     '    return "generator_node"                               # 불합격 -> 다시 생성!\n'
     '\n'
     '# 그래프에 등록\n'
     'workflow.add_conditional_edges("evaluator", route_function, ["generator_node", END])', size=14)
ml(s, Inches(0.8), Inches(5.8), Inches(11.5), Inches(1.2), [
    "핵심: 라우팅 함수가 이전 노드 이름을 반환하면 루프가 됩니다!",
    "무한 루프 방지를 위해 반드시 최대 횟수 제한을 넣어주세요.",
], size=18, spacing=1.5)


# ================================================================
# 슬라이드 6: Evaluator-Optimizer 섹션
# ================================================================
section_slide("PATTERN 1", "Evaluator-Optimizer",
              "생성과 평가를 반복하여 점점 더 나은 결과물을 만드는 루프 패턴")


# ================================================================
# 슬라이드 7: Evaluator-Optimizer 개념
# ================================================================
s = new_slide()
title_line(s, "Evaluator-Optimizer란?", "7")
ml(s, Inches(0.8), Inches(1.4), Inches(11.5), Inches(2.0), [
    "두 개의 LLM이 협력하여 결과물을 반복 개선하는 패턴",
    "",
    "  Generator (생성자) : 결과물을 생성합니다 (예: 광고 문구 작성)",
    "  Evaluator (평가자) : 기준에 따라 평가하고 피드백을 제공합니다",
], size=20, spacing=1.4)
ml(s, Inches(0.8), Inches(3.5), Inches(11.5), Inches(1.0), [
    "  [START] -> [Generator] -> [Evaluator] -> PASS -> [END]",
    "                   ^____________| FAIL (피드백 반영 후 재시도)",
], size=20, bold=True, spacing=1.3, font="Consolas")
ml(s, Inches(0.8), Inches(5.0), Inches(11.5), Inches(2.0), [
    "비유: 신입 카피라이터 + 깐깐한 마케팅 팀장",
    "  신입이 광고 문구 작성 -> 팀장이 검수 -> 불합격이면 피드백 주고 다시 작성!",
    "",
    "이번 예제: 인스타그램 광고 문구를 자동 생성하고 품질 기준에 맞을 때까지 반복 개선",
], size=18, spacing=1.4)


# ================================================================
# 슬라이드 8: AdState 정의
# ================================================================
s = new_slide()
title_line(s, "핵심 코드 1: State 정의 (AdState)", "8")
ml(s, Inches(0.8), Inches(1.4), Inches(11.5), Inches(0.8), [
    "워크플로우에서 사용할 모든 데이터를 미리 선언합니다.",
    "Generator와 Evaluator 모두 이 State를 읽고 씁니다.",
], size=18, spacing=1.5)
code(s, Inches(0.8), Inches(2.6), Inches(11.5), Inches(2.8),
     'class AdState(TypedDict):\n'
     '    product_name: str      # 입력: 광고할 상품명 (예: "자율주행 자동차")\n'
     '    ad_copy: str           # Generator가 작성한 광고 문구\n'
     '    feedback: str          # Evaluator가 제공하는 피드백 (수정 지시사항)\n'
     '    status: str            # 평가 결과: "pass" (합격) / "fail" (불합격)\n'
     '    iteration_count: int   # 현재까지 시도한 횟수 (무한 루프 방지용)', size=16)
ml(s, Inches(0.8), Inches(5.8), Inches(11.5), Inches(1.2), [
    "핵심 포인트:",
    "  - feedback: 불합격 시 \"어떻게 수정하라\"는 구체적 지시사항",
    "  - iteration_count: 무한 루프를 방지하는 안전장치 (최대 3회)",
], size=17, spacing=1.4)


# ================================================================
# 슬라이드 9: EvaluationResult 스키마
# ================================================================
s = new_slide()
title_line(s, "핵심 코드 2: EvaluationResult (평가 출력 형식)", "9")
ml(s, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.8), [
    "Evaluator(팀장)의 응답을 정해진 형식으로 강제합니다.",
    "\"pass/fail\"과 \"피드백\"을 정확히 받아야 루프를 제어할 수 있습니다.",
], size=18, spacing=1.5)
code(s, Inches(0.8), Inches(2.5), Inches(11.5), Inches(2.5),
     'class EvaluationResult(BaseModel):\n'
     '    # Literal: "pass" 또는 "fail" 두 값 중 하나만 허용\n'
     '    status: Literal["pass", "fail"] = Field(\n'
     '        description="기준 충족 여부"\n'
     '    )\n'
     '    # 불합격(fail) 시 구체적인 수정 지시사항\n'
     '    feedback: str = Field(\n'
     '        description="탈락 시 구체적인 수정 지시사항"\n'
     '    )', size=15)
code(s, Inches(0.8), Inches(5.2), Inches(11.5), Inches(1.0),
     '# Evaluator 전용 LLM 생성 (응답이 EvaluationResult 객체로 자동 파싱됨)\n'
     'evaluator_llm = model.with_structured_output(EvaluationResult)', size=15)
ml(s, Inches(0.8), Inches(6.4), Inches(11.5), Inches(0.7), [
    "이렇게 하면 evaluator_llm.invoke(prompt)의 결과가 result.status, result.feedback로 접근 가능!",
], size=16, spacing=1.3)


# ================================================================
# 슬라이드 10: Generator 노드 (copywriter_node)
# ================================================================
s = new_slide()
title_line(s, "핵심 코드 3: Generator 노드 (copywriter_node)", "10")
tb(s, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.4),
   "신입 카피라이터 역할: 광고 문구를 생성합니다", size=18)

tb(s, Inches(0.8), Inches(1.8), Inches(5.7), Inches(0.3),
   "첫 번째 시도 (피드백 없을 때)", size=16, bold=True)
code(s, Inches(0.8), Inches(2.2), Inches(5.7), Inches(2.3),
     'def copywriter_node(state: AdState):\n'
     '    product = state["product_name"]\n'
     '    feedback = state.get("feedback")\n'
     '    count = state.get("iteration_count", 0)\n'
     '\n'
     '    if not feedback:  # 첫 시도\n'
     '        prompt = f"\'{product}\'의 홍보 문구를\n'
     '                   건조하게 작성해줘."', size=12)

tb(s, Inches(6.8), Inches(1.8), Inches(5.7), Inches(0.3),
   "재시도 (피드백 반영)", size=16, bold=True)
code(s, Inches(6.8), Inches(2.2), Inches(5.7), Inches(2.3),
     '    else:  # 피드백이 있으면 반영\n'
     '        prompt = f"""\n'
     "        '{product}' 홍보 문구를 다시 작성해.\n"
     '        <반드시 지켜야 할 수정 사항>\n'
     '        {feedback}\n'
     '        </반드시 지켜야 할 수정 사항>\n'
     '        """', size=12)

ml(s, Inches(0.8), Inches(4.8), Inches(11.5), Inches(2.2), [
    "반환값:",
], size=17, spacing=1.3)
code(s, Inches(0.8), Inches(5.3), Inches(11.5), Inches(1.2),
     '    msg = model.invoke(prompt)\n'
     '    return {"ad_copy": msg.content, "iteration_count": count + 1}\n'
     '    # ad_copy: 새 광고 문구, iteration_count: 시도 횟수 +1', size=14)
ml(s, Inches(0.8), Inches(6.7), Inches(11.5), Inches(0.5), [
    "핵심: 피드백 유무에 따라 프롬프트가 달라진다! 피드백이 있으면 수정 지시를 포함.",
], size=16, spacing=1.3)


# ================================================================
# 슬라이드 11: Evaluator 노드 (manager_node)
# ================================================================
s = new_slide()
title_line(s, "핵심 코드 4: Evaluator 노드 (manager_node)", "11")
tb(s, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.4),
   "깐깐한 마케팅 팀장 역할: 광고 문구를 기준에 따라 심사합니다", size=18)
code(s, Inches(0.8), Inches(1.9), Inches(11.5), Inches(4.0),
     'def manager_node(state: AdState):\n'
     '    ad_copy = state["ad_copy"]   # 신입이 작성한 광고 문구\n'
     '\n'
     '    prompt = f"""\n'
     '    당신은 깐깐한 마케팅 팀장입니다. 다음 광고 문구를 평가하세요:\n'
     '    "{ad_copy}"\n'
     '\n'
     '    <평가 기준>\n'
     '    1. (정량) 해시태그(#)가 3개 이상\n'
     '    2. (정량) \'할인\' 또는 \'특가\' 단어 포함\n'
     '    3. (정성) 활기차고 매력적인 톤 (건조하면 불합격)\n'
     '    </평가 기준>\n'
     '    """\n'
     '\n'
     '    result = evaluator_llm.invoke(prompt)  # EvaluationResult 객체 반환\n'
     '    return {"status": result.status, "feedback": result.feedback}', size=13)
ml(s, Inches(0.8), Inches(6.2), Inches(11.5), Inches(1.0), [
    "핵심: evaluator_llm은 with_structured_output으로 생성 -> result.status, result.feedback 접근 가능",
    "  평가 기준을 프롬프트에 명시적으로 작성할수록 일관된 판단을 얻을 수 있습니다",
], size=16, spacing=1.4)


# ================================================================
# 슬라이드 12: 라우팅 함수 & 그래프 조립
# ================================================================
s = new_slide()
title_line(s, "핵심 코드 5: 라우팅 함수 & 그래프 조립", "12")
tb(s, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.3),
   "라우팅 함수 (루프 제어)", size=16, bold=True)
code(s, Inches(0.8), Inches(1.7), Inches(5.7), Inches(2.8),
     'def route_submission(state: AdState):\n'
     '    status = state["status"]\n'
     '    count = state["iteration_count"]\n'
     '\n'
     '    # 합격 -> 종료\n'
     '    if status == "pass":\n'
     '        return END\n'
     '    # 3번 시도 초과 -> 강제 종료\n'
     '    if count >= 3:\n'
     '        return END\n'
     '    # 불합격 -> 다시 Generator로!\n'
     '    return "copywriter_node"', size=13)

tb(s, Inches(6.8), Inches(1.3), Inches(5.5), Inches(0.3),
   "그래프 조립", size=16, bold=True)
code(s, Inches(6.8), Inches(1.7), Inches(5.7), Inches(2.8),
     'workflow = StateGraph(AdState)\n'
     '\n'
     '# 노드 등록\n'
     'workflow.add_node("copywriter_node",\n'
     '                  copywriter_node)\n'
     'workflow.add_node("manager_node",\n'
     '                  manager_node)\n'
     '\n'
     '# 엣지 연결\n'
     'workflow.add_edge(START, "copywriter_node")\n'
     'workflow.add_edge("copywriter_node",\n'
     '                  "manager_node")', size=13)

code(s, Inches(0.8), Inches(4.8), Inches(11.5), Inches(1.5),
     '# 조건부 엣지: manager_node 실행 후 route_submission 결과에 따라 분기\n'
     'workflow.add_conditional_edges(\n'
     '    "manager_node",           # 이 노드 실행 후\n'
     '    route_submission,         # 이 함수가 다음 노드를 결정\n'
     '    ["copywriter_node", END]  # 가능한 목적지: 다시 생성 or 종료\n'
     ')\n'
     'app = workflow.compile()', size=14)
ml(s, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.7), [
    "핵심: manager_node -> route_submission -> copywriter_node (루프) 또는 END (종료)",
], size=17, spacing=1.3)


# ================================================================
# 슬라이드 13: 실행 & 정리
# ================================================================
s = new_slide()
title_line(s, "실행 & 정리: Evaluator-Optimizer", "13")
code(s, Inches(0.8), Inches(1.4), Inches(11.5), Inches(2.0),
     '# 실행: 상품명만 입력하면 됩니다!\n'
     'result = app.invoke({"product_name": "자율주행 자동차"})\n'
     '\n'
     '# 결과 확인\n'
     'print(result["ad_copy"])          # 최종 광고 문구\n'
     'print(result["status"])           # pass or fail\n'
     'print(result["iteration_count"])  # 몇 번 시도했는지', size=15)
ml(s, Inches(0.8), Inches(3.8), Inches(11.5), Inches(3.2), [
    "Evaluator-Optimizer 핵심 정리:",
    "",
    "  구조: Generator -> Evaluator -> (합격이면 END, 불합격이면 다시 Generator)",
    "  핵심 기술: Conditional Edge로 루프 구현 + Structured Output으로 평가 결과 파싱",
    "  안전장치: iteration_count로 최대 시도 횟수 제한 (무한 루프 방지)",
    "",
    "  사용처: 콘텐츠 품질 자동 검수, 코드 리뷰 자동화, 보고서 개선 루프 등",
    "  비유: 신입이 수정 -> 팀장이 검수 -> 부족하면 다시 수정 -> 합격할 때까지 반복!",
], size=18, spacing=1.3)


# ================================================================
# 슬라이드 14: Orchestrator-Worker 섹션
# ================================================================
section_slide("PATTERN 2", "Orchestrator-Worker",
              "큰 작업을 분할하여 병렬 처리하는 Map-Reduce 패턴")


# ================================================================
# 슬라이드 15: Orchestrator-Worker 개념
# ================================================================
s = new_slide()
title_line(s, "Orchestrator-Worker란?", "15")
ml(s, Inches(0.8), Inches(1.4), Inches(11.5), Inches(1.5), [
    "Orchestrator (지휘자): 큰 작업을 여러 하위 작업으로 분배",
    "Worker (작업자): 각자 맡은 작업을 독립적으로 수행 (병렬 실행)",
    "Synthesizer (통합자): 모든 결과를 하나로 합침",
], size=20, spacing=1.5)
ml(s, Inches(0.8), Inches(3.2), Inches(11.5), Inches(2.0), [
    "                   +-> [Worker 1] --+",
    "START -> [Orchestrator] -> [Worker 2] --> [Synthesizer] -> END",
    "                   +-> [Worker 3] --+",
    "         (계획 수립)    (병렬 실행)         (결과 취합)",
], size=18, spacing=1.2, font="Consolas")
ml(s, Inches(0.8), Inches(5.3), Inches(11.5), Inches(1.7), [
    "비유: 편집장(Orchestrator)이 기자들(Worker)에게 각 섹션을 분배 -> 동시 집필 -> 편집자(Synthesizer)가 합본",
    "",
    "이번 예제: 주제를 주면 LLM이 목차를 계획하고, 각 섹션을 병렬로 작성 후 하나의 보고서로 합침",
    "핵심 기술: LangGraph의 Send API로 동적으로 Worker를 생성!",
], size=18, spacing=1.4)


# ================================================================
# 슬라이드 16: ReportState 정의
# ================================================================
s = new_slide()
title_line(s, "핵심 코드 1: State 정의 (ReportState)", "16")
ml(s, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.8), [
    "핵심 개념: Annotated + operator.add",
    "여러 Worker가 동시에 결과를 반환할 때, 덮어쓰지 않고 누적해야 합니다!",
], size=18, bold=True, spacing=1.5)
code(s, Inches(0.8), Inches(2.5), Inches(11.5), Inches(3.0),
     'import operator\n'
     'from typing import Annotated, List\n'
     '\n'
     'class ReportState(TypedDict):\n'
     '    topic: str                    # 보고서 주제 (입력값)\n'
     '    sections: List[Section]       # Orchestrator가 계획한 목차 목록\n'
     '\n'
     '    # Annotated + operator.add = 여러 Worker 결과를 자동 누적!\n'
     '    # Worker1 -> ["섹션1"], Worker2 -> ["섹션2"]\n'
     '    # => completed_sections = ["섹션1", "섹션2"]  (덮어쓰기 X, 합치기 O)\n'
     '    completed_sections: Annotated[List[str], operator.add]\n'
     '\n'
     '    final_report: str             # Synthesizer가 완성한 최종 보고서', size=14)
ml(s, Inches(0.8), Inches(5.8), Inches(11.5), Inches(1.2), [
    "핵심: Annotated[List[str], operator.add]가 없으면 나중에 반환한 Worker의 결과만 남음!",
    "  이 선언 덕분에 여러 Worker의 결과가 충돌 없이 하나의 리스트로 합쳐집니다.",
], size=17, spacing=1.4)


# ================================================================
# 슬라이드 17: Section & Plan 스키마
# ================================================================
s = new_slide()
title_line(s, "핵심 코드 2: Section & Plan (Pydantic 스키마)", "17")
ml(s, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.5), [
    "Orchestrator가 LLM으로 목차를 생성할 때 사용하는 출력 형식입니다.",
], size=18, spacing=1.5)

tb(s, Inches(0.8), Inches(2.0), Inches(5.5), Inches(0.3),
   "Section: 각 목차 항목", size=16, bold=True)
code(s, Inches(0.8), Inches(2.4), Inches(5.7), Inches(1.8),
     'class Section(BaseModel):\n'
     '    name: str = Field(\n'
     '        description="목차의 제목"\n'
     '    )\n'
     '    description: str = Field(\n'
     '        description="핵심 내용 가이드"\n'
     '    )', size=14)

tb(s, Inches(6.8), Inches(2.0), Inches(5.5), Inches(0.3),
   "Plan: 전체 계획표", size=16, bold=True)
code(s, Inches(6.8), Inches(2.4), Inches(5.7), Inches(1.8),
     'class Plan(BaseModel):\n'
     '    sections: List[Section] = Field(\n'
     '        description="보고서 목차 리스트"\n'
     '    )\n'
     '\n'
     'planner_llm = model.with_structured_output(Plan)', size=14)

tb(s, Inches(0.8), Inches(4.5), Inches(11.5), Inches(0.3),
   "WorkerState: Worker 전용 상태 (작업 지시서)", size=16, bold=True)
code(s, Inches(0.8), Inches(4.9), Inches(11.5), Inches(1.2),
     '# Worker는 전체 ReportState가 아닌, 자기가 맡은 Section만 알면 됩니다\n'
     'class WorkerState(TypedDict):\n'
     '    section: Section  # 이 Worker가 작성해야 할 단일 섹션 정보', size=14)

ml(s, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.8), [
    "핵심: Section으로 목차 구조를 잡고, Plan으로 전체 계획을 받고, WorkerState로 개별 작업을 전달!",
], size=17, spacing=1.3)


# ================================================================
# 슬라이드 18: Orchestrator 노드
# ================================================================
s = new_slide()
title_line(s, "핵심 코드 3: Orchestrator 노드 (목차 계획)", "18")
ml(s, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.5), [
    "주어진 주제에 대해 보고서 구조(목차)를 설계합니다.",
], size=18, spacing=1.5)
code(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(3.0),
     'def orchestrator_node(state: ReportState):\n'
     '    topic = state["topic"]\n'
     '\n'
     '    # Structured Output으로 목차를 Plan 객체로 받아옴\n'
     '    plan = planner_llm.invoke(\n'
     '        f"\'{topic}\'에 대한 보고서 목차를 짜줘. 3개 섹션 이내로 구성해."\n'
     '    )\n'
     '    # plan.sections = [Section(name=..., description=...), ...]\n'
     '\n'
     '    return {"sections": plan.sections}  # 목차를 State에 저장', size=15)
ml(s, Inches(0.8), Inches(5.3), Inches(11.5), Inches(1.7), [
    "실행 흐름:",
    "  1. topic(\"생성형 AI의 미래\")을 읽음",
    "  2. planner_llm이 구조화된 목차(Plan)를 생성",
    "  3. sections 리스트를 State에 저장 -> 다음 단계에서 Worker들에게 분배됨",
], size=17, spacing=1.4)


# ================================================================
# 슬라이드 19: Worker 노드
# ================================================================
s = new_slide()
title_line(s, "핵심 코드 4: Worker 노드 (섹션 집필)", "19")
code(s, Inches(0.8), Inches(1.4), Inches(11.5), Inches(3.5),
     'def worker_node(state: WorkerState):\n'
     '    section = state["section"]  # 이 Worker에게 할당된 섹션\n'
     '\n'
     '    prompt = f"""\n'
     '    다음 섹션에 대한 내용을 짧게 작성해줘.\n'
     '    제목: {section.name}\n'
     '    내용 가이드: {section.description}\n'
     '    """\n'
     '    msg = model.invoke(prompt)\n'
     '\n'
     '    content = f"## {section.name}\\n{msg.content}\\n"\n'
     '\n'
     '    # 리스트로 감싸서 반환! (operator.add로 누적되기 때문)\n'
     '    return {"completed_sections": [content]}', size=15)
ml(s, Inches(0.8), Inches(5.2), Inches(11.5), Inches(2.0), [
    "핵심 포인트:",
    "  - 입력 State는 WorkerState (전체 ReportState가 아닌 개별 Section만!)",
    "  - 반환값을 리스트 [content]로 감싸야 operator.add가 누적할 수 있음",
    "  - 여러 Worker가 동시에 실행되어도 각자의 결과가 안전하게 합쳐짐",
], size=17, spacing=1.4)


# ================================================================
# 슬라이드 20: Send API & 그래프 조립
# ================================================================
s = new_slide()
title_line(s, "핵심 코드 5: Send API & 그래프 조립", "20")
ml(s, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.5), [
    "Send API: 런타임에 동적으로 여러 Worker를 생성하여 병렬 실행!",
], size=18, bold=True, spacing=1.5)

tb(s, Inches(0.8), Inches(2.0), Inches(5.5), Inches(0.3),
   "동적 라우팅 함수 (Map 단계)", size=16, bold=True)
code(s, Inches(0.8), Inches(2.4), Inches(5.7), Inches(2.2),
     'from langgraph.types import Send\n'
     '\n'
     'def assign_workers(state: ReportState):\n'
     '    sections = state["sections"]\n'
     '    # 섹션 수만큼 Worker를 동적 생성!\n'
     '    return [\n'
     '        Send("worker_node", {"section": s})\n'
     '        for s in sections\n'
     '    ]', size=13)

tb(s, Inches(6.8), Inches(2.0), Inches(5.5), Inches(0.3),
   "그래프 조립", size=16, bold=True)
code(s, Inches(6.8), Inches(2.4), Inches(5.7), Inches(2.2),
     'workflow = StateGraph(ReportState)\n'
     'workflow.add_node("orchestrator_node", ...)\n'
     'workflow.add_node("worker_node", ...)\n'
     'workflow.add_node("synthesizer_node", ...)\n'
     '\n'
     'workflow.add_edge(START, "orchestrator_node")\n'
     'workflow.add_edge("worker_node",\n'
     '                  "synthesizer_node")\n'
     'workflow.add_edge("synthesizer_node", END)', size=13)

code(s, Inches(0.8), Inches(4.9), Inches(11.5), Inches(1.5),
     '# 핵심: Orchestrator -> (Send API로 Map) -> Worker들 (병렬)\n'
     'workflow.add_conditional_edges(\n'
     '    "orchestrator_node",\n'
     '    assign_workers,        # Send 리스트를 반환하는 함수\n'
     '    ["worker_node"]        # 가능한 목적지\n'
     ')\n'
     'app = workflow.compile()', size=14)
ml(s, Inches(0.8), Inches(6.6), Inches(11.5), Inches(0.6), [
    "핵심: Send(\"worker_node\", {state}) -> 섹션 3개면 Worker 3개가 동시에 실행!",
], size=17, spacing=1.3)


# ================================================================
# 슬라이드 21: Synthesizer v1 vs v2
# ================================================================
s = new_slide()
title_line(s, "핵심 코드 6: Synthesizer (결과 취합)", "21")

tb(s, Inches(0.8), Inches(1.3), Inches(5.7), Inches(0.3),
   "v1: 단순 합치기", size=16, bold=True)
code(s, Inches(0.8), Inches(1.7), Inches(5.7), Inches(2.0),
     'def synthesizer_node(state: ReportState):\n'
     '    completed = state["completed_sections"]\n'
     '    # 단순히 줄바꿈으로 연결\n'
     '    final = "\\n".join(completed)\n'
     '    return {"final_report": final}', size=13)

tb(s, Inches(6.8), Inches(1.3), Inches(5.7), Inches(0.3),
   "v2: LLM 편집자가 재구성", size=16, bold=True)
code(s, Inches(6.8), Inches(1.7), Inches(5.7), Inches(2.0),
     'def synthesizer_node(state: ReportState):\n'
     '    raw = "\\n\\n".join(\n'
     '        state["completed_sections"])\n'
     '    prompt = f"전문 편집자로서 하나의\n'
     '    자연스러운 보고서로 재구성해줘..."\n'
     '    msg = model.invoke(prompt)\n'
     '    return {"final_report": msg.content}', size=13)

ml(s, Inches(0.8), Inches(4.2), Inches(11.5), Inches(2.8), [
    "v1 vs v2 비교:",
    "",
    "  v1 (단순 합치기): 각 Worker의 결과를 그대로 이어붙임 -> 빠르지만 어색한 연결",
    "  v2 (LLM 편집):  LLM 편집자가 전체 내용을 재구성 -> 서론/결론 추가, 매끄러운 연결",
    "",
    "  v2의 프롬프트에 \"서론과 결론 추가\", \"마크다운 형식\" 등 구체적 지시를 포함!",
    "  같은 그래프 구조에서 synthesizer_node 함수만 교체하면 업그레이드 완료!",
], size=17, spacing=1.3)


# ================================================================
# 슬라이드 22: Orchestrator-Worker 실행 & 정리
# ================================================================
s = new_slide()
title_line(s, "실행 & 정리: Orchestrator-Worker", "22")
code(s, Inches(0.8), Inches(1.4), Inches(11.5), Inches(1.6),
     '# 실행: 주제만 입력하면 보고서가 자동 생성됩니다!\n'
     'result = app.invoke({"topic": "생성형 AI의 미래"})\n'
     'print(result["final_report"])', size=15)
ml(s, Inches(0.8), Inches(3.3), Inches(11.5), Inches(3.7), [
    "Orchestrator-Worker 핵심 정리:",
    "",
    "  구조: Orchestrator(계획) -> Worker들(병렬 집필) -> Synthesizer(취합)",
    "  핵심 기술:",
    "    - Send API: 런타임에 동적으로 Worker 수를 결정하여 병렬 실행",
    "    - Annotated + operator.add: 여러 Worker의 결과를 안전하게 누적",
    "    - Structured Output: 목차를 구조화된 형식으로 생성",
    "",
    "  사용처: 보고서 자동 작성, 대규모 데이터 분석, 다중 관점 리서치 등",
    "  비유: 편집장이 기자들에게 각 섹션 분배 -> 동시 집필 -> 편집자가 합본!",
], size=18, spacing=1.2)


# ================================================================
# 슬라이드 23: Agents 섹션
# ================================================================
section_slide("PATTERN 3", "Agents (에이전트)",
              "LLM이 스스로 판단하여 도구(Tool)를 선택하고 실행하는 자율 패턴")


# ================================================================
# 슬라이드 24: Agents 개념
# ================================================================
s = new_slide()
title_line(s, "Agents (에이전트)란?", "24")
ml(s, Inches(0.8), Inches(1.4), Inches(11.5), Inches(1.5), [
    "LLM이 스스로 판단하여 도구(Tool)를 선택하고 실행하는 패턴",
    "단순히 텍스트를 생성하는 것을 넘어, 외부 함수를 직접 호출할 수 있습니다!",
], size=20, spacing=1.5)
ml(s, Inches(0.8), Inches(3.0), Inches(11.5), Inches(2.0), [
    "동작 원리 (ReAct 패턴):",
    "",
    "  [START] -> [LLM 판단] -> 도구 필요? -> YES -> [도구 실행] -> [LLM 판단] (반복)",
    "                                      -> NO  -> [최종 답변] -> [END]",
], size=18, spacing=1.3, font="Consolas")
ml(s, Inches(0.8), Inches(5.2), Inches(11.5), Inches(1.8), [
    "비유: 만능 비서에게 \"3+4를 구하고, 그 결과에 7을 곱해줘\" 라고 말하면",
    "  비서가 알아서 계산기(도구)를 꺼내서 3+4=7 계산 -> 다시 7*7=49 계산 -> 최종 답변",
    "",
    "이번 예제: 덧셈, 곱셈, 나눗셈 도구를 LLM에게 주고, 복합 계산을 자율적으로 처리",
], size=18, spacing=1.4)


# ================================================================
# 슬라이드 25: Tool 정의 & bind_tools
# ================================================================
s = new_slide()
title_line(s, "핵심 코드 1: Tool 정의 & bind_tools", "25")
tb(s, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.3),
   "@tool 데코레이터로 함수를 도구로 변환", size=16, bold=True)
code(s, Inches(0.8), Inches(1.7), Inches(5.7), Inches(3.0),
     'from langchain.tools import tool\n'
     '\n'
     '@tool\n'
     'def multiply(a: int, b: int) -> int:\n'
     '    """Multiply `a` and `b`.\n'
     '    Args:\n'
     '        a: First int\n'
     '        b: Second int\n'
     '    """\n'
     '    return a * b\n'
     '\n'
     '# add, divide도 동일한 패턴!', size=13)

tb(s, Inches(6.8), Inches(1.3), Inches(5.5), Inches(0.3),
   "LLM에 도구 바인딩", size=16, bold=True)
code(s, Inches(6.8), Inches(1.7), Inches(5.7), Inches(3.0),
     '# 도구 목록 정의\n'
     'tools = [add, multiply, divide]\n'
     '\n'
     '# 이름으로 빠르게 찾기 위한 딕셔너리\n'
     'tools_by_name = {\n'
     '    tool.name: tool for tool in tools\n'
     '}\n'
     '\n'
     '# LLM에 도구 바인딩\n'
     '# "이 도구들을 사용할 수 있다"고 알려줌\n'
     'llm_with_tools = model.bind_tools(tools)', size=13)

ml(s, Inches(0.8), Inches(5.0), Inches(11.5), Inches(2.0), [
    "핵심 포인트:",
    "  - @tool 데코레이터: 일반 함수를 LangChain Tool로 변환",
    "  - docstring이 Tool의 설명이 됨 -> LLM이 이 설명을 읽고 언제 쓸지 판단!",
    "  - bind_tools(): LLM이 필요에 따라 도구 호출(tool_calls)을 응답에 포함할 수 있게 됨",
], size=17, spacing=1.4)


# ================================================================
# 슬라이드 26: MessagesState & 노드
# ================================================================
s = new_slide()
title_line(s, "핵심 코드 2: LLM 판단 노드 & 도구 실행 노드", "26")

tb(s, Inches(0.8), Inches(1.3), Inches(5.7), Inches(0.3),
   "LLM 판단 노드 (llm_call)", size=16, bold=True)
code(s, Inches(0.8), Inches(1.7), Inches(5.7), Inches(2.5),
     'from langgraph.graph import MessagesState\n'
     '\n'
     'def llm_call(state: MessagesState):\n'
     '    return {"messages": [\n'
     '        llm_with_tools.invoke(\n'
     '            [SystemMessage(content=\n'
     '                "You are a helpful assistant...")]\n'
     '            + state["messages"]\n'
     '        )\n'
     '    ]}', size=12)

tb(s, Inches(6.8), Inches(1.3), Inches(5.7), Inches(0.3),
   "도구 실행 노드 (tool_node)", size=16, bold=True)
code(s, Inches(6.8), Inches(1.7), Inches(5.7), Inches(2.5),
     'def tool_node(state: dict):\n'
     '    result = []\n'
     '    for tc in state["messages"][-1].tool_calls:\n'
     '        tool = tools_by_name[tc["name"]]\n'
     '        obs = tool.invoke(tc["args"])\n'
     '        result.append(ToolMessage(\n'
     '            content=obs,\n'
     '            tool_call_id=tc["id"]\n'
     '        ))\n'
     '    return {"messages": result}', size=12)

ml(s, Inches(0.8), Inches(4.5), Inches(11.5), Inches(2.5), [
    "MessagesState: messages 키를 가진 특수 State (대화 히스토리 자동 관리)",
    "",
    "llm_call: LLM이 다음 행동을 결정",
    "  - 도구가 필요하면 -> tool_calls가 포함된 응답 반환",
    "  - 충분히 답변 가능하면 -> 일반 텍스트 응답 반환",
    "",
    "tool_node: LLM이 요청한 도구를 실제로 실행하고 결과를 ToolMessage로 반환",
], size=16, spacing=1.2)


# ================================================================
# 슬라이드 27: 라우팅 & 그래프 조립
# ================================================================
s = new_slide()
title_line(s, "핵심 코드 3: 라우팅 & 그래프 조립", "27")

tb(s, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.3),
   "라우팅 함수 (루프 제어)", size=16, bold=True)
code(s, Inches(0.8), Inches(1.7), Inches(5.7), Inches(2.3),
     'def should_continue(state: MessagesState):\n'
     '    last = state["messages"][-1]\n'
     '\n'
     '    # 도구 호출 요청이 있으면\n'
     '    # -> tool_node로 이동\n'
     '    if last.tool_calls:\n'
     '        return "tool_node"\n'
     '\n'
     '    # 도구 호출 없으면\n'
     '    # -> 최종 답변 완성, 종료\n'
     '    return END', size=13)

tb(s, Inches(6.8), Inches(1.3), Inches(5.5), Inches(0.3),
   "에이전트 그래프 조립", size=16, bold=True)
code(s, Inches(6.8), Inches(1.7), Inches(5.7), Inches(2.3),
     'agent = StateGraph(MessagesState)\n'
     '\n'
     'agent.add_node("llm_call", llm_call)\n'
     'agent.add_node("tool_node", tool_node)\n'
     '\n'
     'agent.add_edge(START, "llm_call")\n'
     'agent.add_conditional_edges(\n'
     '    "llm_call", should_continue,\n'
     '    ["tool_node", END]\n'
     ')\n'
     'agent.add_edge("tool_node", "llm_call")', size=13)

ml(s, Inches(0.8), Inches(4.3), Inches(11.5), Inches(2.7), [
    "루프 구조:",
    "  START -> llm_call -> (도구 필요?) -> tool_node -> llm_call -> ... -> END",
    "",
    "핵심:",
    "  - llm_call -> tool_node -> llm_call 이 루프가 에이전트의 심장!",
    "  - LLM이 \"더 이상 도구가 필요 없다\"고 판단하면 자동으로 종료",
    "  - 복잡한 작업도 여러 번의 도구 호출로 단계적으로 처리 가능",
], size=17, spacing=1.3)


# ================================================================
# 슬라이드 28: 실행 & 메시지 흐름
# ================================================================
s = new_slide()
title_line(s, "실행 & 정리: Agents", "28")
code(s, Inches(0.8), Inches(1.4), Inches(11.5), Inches(2.2),
     '# 복합 요청: 덧셈 -> 곱셈 -> 이야기 생성\n'
     'messages = [HumanMessage(\n'
     '    content="Add 3 and 4. Multiply the output by 7. And make a funny story about Japan."\n'
     ')]\n'
     'result = agent.invoke({"messages": messages})', size=14)
ml(s, Inches(0.8), Inches(3.9), Inches(11.5), Inches(3.2), [
    "실행 흐름 (메시지 히스토리):",
    "",
    "  1. HumanMessage: \"Add 3 and 4. Multiply by 7. Make a story.\"",
    "  2. AIMessage (tool_calls): add(3, 4) 호출 요청",
    "  3. ToolMessage: 결과 = 7",
    "  4. AIMessage (tool_calls): multiply(7, 7) 호출 요청",
    "  5. ToolMessage: 결과 = 49",
    "  6. AIMessage: 최종 답변 (계산 결과 + 일본 이야기)",
], size=17, spacing=1.2)


# ================================================================
# 슬라이드 29: 3가지 패턴 비교표
# ================================================================
s = new_slide()
title_line(s, "3가지 패턴 비교표", "29")
ml(s, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5), [
    "+----------------+---------------------+---------------------+---------------------+",
    "|                | Evaluator-Optimizer | Orchestrator-Worker | Agents              |",
    "+----------------+---------------------+---------------------+---------------------+",
    "| 구조           | Generator -> Eval   | Orch -> Workers     | LLM -> Tool         |",
    "|                | -> (루프 or 종료)   | -> Synthesizer      | -> (루프 or 종료)   |",
    "+----------------+---------------------+---------------------+---------------------+",
    "| 핵심 기술      | Conditional Edge    | Send API            | bind_tools          |",
    "|                | Structured Output   | operator.add        | tool_calls          |",
    "+----------------+---------------------+---------------------+---------------------+",
    "| 장점           | 반복 개선으로       | 동적 병렬 처리      | LLM이 자율적으로    |",
    "|                | 품질 점진적 향상    | 대규모 작업 분할    | 도구 선택/실행      |",
    "+----------------+---------------------+---------------------+---------------------+",
    "| 사용 사례      | 콘텐츠 품질 검수    | 보고서 자동 작성    | 계산, 검색, API     |",
    "|                | 코드 리뷰 자동화    | 다중 관점 분석      | 복합 작업 처리      |",
    "+----------------+---------------------+---------------------+---------------------+",
], size=14, spacing=1.0, font="Consolas")


# ================================================================
# 슬라이드 30: 핵심 개념 총정리
# ================================================================
s = new_slide()
title_line(s, "핵심 개념 총정리", "30")
ml(s, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5), [
    "Conditional Edge (조건부 엣지)",
    "  add_conditional_edges()로 동적 분기 & 루프 구현",
    "",
    "Structured Output (구조화 출력)",
    "  Pydantic BaseModel + with_structured_output()로 LLM 응답 형식 강제",
    "",
    "Send API",
    "  런타임에 동적으로 여러 노드를 생성하여 병렬 실행 (Map-Reduce)",
    "",
    "Annotated + operator.add",
    "  여러 Worker의 결과를 덮어쓰지 않고 리스트로 누적",
    "",
    "Tool & bind_tools",
    "  @tool로 함수를 도구로 변환, bind_tools()로 LLM에 연결",
    "",
    "MessagesState",
    "  에이전트 패턴의 대화 히스토리 자동 관리 State",
], size=18, spacing=1.05)


# ================================================================
# 슬라이드 31: Q&A
# ================================================================
s = new_slide()
tb(s, Inches(1), Inches(2.0), Inches(11), Inches(1.2),
   "Q & A", size=60, bold=True, align=PP_ALIGN.CENTER, font="Consolas")
tb(s, Inches(1), Inches(3.5), Inches(11), Inches(0.8),
   "질문이 있으신가요?", size=28, align=PP_ALIGN.CENTER)
ml(s, Inches(1), Inches(5.0), Inches(11), Inches(1.5), [
    "오늘 배운 3가지 패턴을 직접 실습해보세요!",
    "주제와 평가 기준을 바꿔가며 코드를 실행해보면 이해가 더 깊어집니다.",
], size=20, spacing=1.5, align=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────
# 저장
# ──────────────────────────────────────────────
output_path = r"c:\DL_Practice\[3강]_랭그래프_응용_강의슬라이드_second.pptx"
prs.save(output_path)
print(f"PPT 생성 완료: {output_path}")
print(f"총 슬라이드 수: {len(prs.slides)}")
