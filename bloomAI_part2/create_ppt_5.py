# -*- coding: utf-8 -*-
"""
LangGraph 응용 5강 - Multi-Agent 시스템 강의 슬라이드 자동 생성
(흰 배경 + 검은 글씨 템플릿)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ──────────────────────────────────────────────
# 색상 (흰 배경 / 검은 글씨 / 코드 블록만 연회색)
# ──────────────────────────────────────────────
BLACK      = RGBColor(0x00, 0x00, 0x00)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_LIGHT = RGBColor(0xF0, 0xF0, 0xF0)  # 코드 블록 배경
GRAY_MID   = RGBColor(0x99, 0x99, 0x99)  # 보조 텍스트

# ──────────────────────────────────────────────
# 프레젠테이션 초기화 (16:9)
# ──────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height


# ──────────────────────────────────────────────
# 유틸리티
# ──────────────────────────────────────────────
def new_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def tb(slide, left, top, width, height, text, size=18, bold=False,
       align=PP_ALIGN.LEFT, color=BLACK, font="맑은 고딕", anchor=MSO_ANCHOR.TOP):
    """텍스트 박스 추가 (단일 텍스트)"""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    try:
        tf.vertical_anchor = anchor
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return box


def ml(slide, left, top, width, height, lines, size=16, bold=False,
       align=PP_ALIGN.LEFT, color=BLACK, spacing=1.5, font="맑은 고딕"):
    """멀티라인 텍스트 박스"""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font
        p.alignment = align
        p.space_after = Pt(size * (spacing - 1))
    return box


def code(slide, left, top, width, height, text, size=13):
    """코드 블록 (연회색 배경 + Consolas)"""
    add_rect(slide, left, top, width, height, GRAY_LIGHT)
    box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15),
                                   width - Inches(0.4), height - Inches(0.3))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = BLACK
        p.font.name = "Consolas"
        p.space_after = Pt(2)


def title_line(slide, text, num=""):
    """슬라이드 상단 제목 + 밑줄"""
    tb(slide, Inches(0.6), Inches(0.3), Inches(11), Inches(0.7),
       text, size=28, bold=True)
    add_rect(slide, Inches(0.6), Inches(1.05), Inches(12), Pt(2), BLACK)
    if num:
        tb(slide, Inches(12), Inches(0.3), Inches(0.8), Inches(0.7),
           num, size=14, color=GRAY_MID, align=PP_ALIGN.RIGHT)


def section_slide(number, title, subtitle=""):
    """섹션 구분 슬라이드"""
    s = new_slide()
    tb(s, Inches(1), Inches(1.8), Inches(11), Inches(1),
       number, size=48, bold=True, align=PP_ALIGN.CENTER, font="Consolas")
    tb(s, Inches(1), Inches(3.2), Inches(11), Inches(1),
       title, size=36, bold=True, align=PP_ALIGN.CENTER)
    if subtitle:
        tb(s, Inches(1), Inches(4.5), Inches(11), Inches(0.7),
           subtitle, size=20, color=GRAY_MID, align=PP_ALIGN.CENTER)
    return s


# ================================================================
# 슬라이드 1: 표지
# ================================================================
s = new_slide()
tb(s, Inches(1), Inches(1.5), Inches(11), Inches(1),
   "LangGraph 응용 강의", size=48, bold=True, align=PP_ALIGN.CENTER)
tb(s, Inches(1), Inches(2.8), Inches(11), Inches(0.8),
   "5강: Multi-Agent 시스템", size=30, bold=True, align=PP_ALIGN.CENTER)
tb(s, Inches(1), Inches(4.2), Inches(11), Inches(0.6),
   "Interrupt  |  Command  |  Subgraph  |  Supervisor  |  Handoff",
   size=22, align=PP_ALIGN.CENTER)
tb(s, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
   "초보자를 위한 단계별 실습 가이드", size=18, color=GRAY_MID, align=PP_ALIGN.CENTER)


# ================================================================
# 슬라이드 2: 학습 목표
# ================================================================
s = new_slide()
title_line(s, "학습 목표", "2")
ml(s, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.0), [
    "이번 강의에서 배울 8가지 핵심 패턴:",
    "",
    "1. Dynamic Interrupt       — 실행 중 사용자 입력을 받고 재개하기",
    "2. Interrupts in Tools     — 도구(Tool) 안에서 사용자에게 질문하기",
    "3. Validating Human Input  — 사용자 입력을 검증하는 루프 만들기",
    "4. Command with Goto       — 런타임에 동적으로 노드 이동하기",
    "5. Subgraph (State 공유)   — 부모-자식 그래프가 State를 공유",
    "6. Subgraph (State 분리)   — 부모-자식 그래프가 각자 State 사용",
    "7. Multi-Agent: Supervisor — 중앙 관리자가 워커를 지휘하는 패턴",
    "8. Multi-Agent: Handoff    — 에이전트끼리 직접 업무를 넘기는 패턴",
], size=19, spacing=1.3)


# ================================================================
# 슬라이드 3: 사전 지식 복습
# ================================================================
s = new_slide()
title_line(s, "사전 지식 복습: LangGraph 핵심 개념", "3")
ml(s, Inches(0.8), Inches(1.4), Inches(11.5), Inches(5.0), [
    "State (상태)",
    "  모든 노드가 공유하는 데이터 저장소 (TypedDict로 정의)",
    "  = 칠판에 적힌 메모 (모두가 읽고 쓸 수 있음)",
    "",
    "Node (노드)",
    "  State를 처리하는 함수 (입력: State, 출력: 업데이트할 딕셔너리)",
    "  = 각 단계의 작업자 (한 가지 일을 담당)",
    "",
    "Edge (엣지)",
    "  노드 간의 연결 (실행 순서를 정하는 화살표)",
    "  add_edge(A, B) = 고정 연결  |  add_conditional_edges() = 조건부 연결",
    "",
    "Compile & Invoke",
    "  app = workflow.compile()  →  result = app.invoke(입력)",
    "  기본 구조:  START → Node A → Node B → END",
], size=18, spacing=1.2)


# ================================================================
# 슬라이드 4: 이번 강의 로드맵
# ================================================================
s = new_slide()
title_line(s, "이번 강의 로드맵", "4")
ml(s, Inches(0.8), Inches(1.4), Inches(11.5), Inches(1.5), [
    "기초부터 고급까지, 단계적으로 배워갑니다:",
], size=20, spacing=1.5)
ml(s, Inches(0.8), Inches(2.3), Inches(11.5), Inches(4.5), [
    "[기초]  1~3: Interrupt 패턴",
    "  실행 중 멈추고 사용자에게 질문 → 답변 받고 이어서 실행",
    "",
    "[중급]  4: Command(goto) 패턴",
    "  런타임에 \"다음에 어디로 갈지\" 동적으로 결정",
    "",
    "[고급]  5~6: Subgraph 패턴",
    "  그래프 안에 그래프를 넣어 모듈화 (레고 블록처럼 조립)",
    "",
    "[심화]  7~8: Multi-Agent 패턴",
    "  여러 AI 에이전트가 협력하여 복잡한 작업 수행",
], size=18, spacing=1.2)


# ================================================================
# 슬라이드 5: SECTION 1 구분
# ================================================================
section_slide("SECTION 1", "Dynamic Interrupt",
              "그래프 실행 중 멈추고 사용자 입력을 받는 패턴")


# ================================================================
# 슬라이드 6: Dynamic Interrupt 이론
# ================================================================
s = new_slide()
title_line(s, "Dynamic Interrupt란?", "6")
ml(s, Inches(0.8), Inches(1.4), Inches(11.5), Inches(2.0), [
    "그래프가 실행 중간에 \"잠깐 멈춤\" → 사용자에게 질문 → 답변 받고 \"이어서 실행\"",
    "",
    "비유: 은행 창구",
    "  직원이 서류를 처리하다가 → \"주소를 확인해주세요\" → 고객이 답변 → 처리 재개",
], size=20, spacing=1.4)
ml(s, Inches(0.8), Inches(3.8), Inches(11.5), Inches(3.0), [
    "왜 필요한가?",
    "  - AI가 작업 중 추가 정보가 필요할 때",
    "  - 중요한 결정을 사용자에게 확인받고 싶을 때",
    "  - 단계별로 사용자 승인이 필요한 워크플로우",
    "",
    "핵심 키워드:",
    "  interrupt()  — 실행을 멈추고 질문",
    "  Command(resume=값)  — 답변을 보내고 실행 재개",
    "  InMemorySaver  — 멈춘 시점의 상태를 저장하는 체크포인터",
], size=18, spacing=1.3)


# ================================================================
# 슬라이드 7: ChatState 정의
# ================================================================
s = new_slide()
title_line(s, "핵심 코드 1: ChatState 정의", "7")
ml(s, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.8), [
    "Interrupt 패턴에서 사용하는 State 구조를 정의합니다.",
], size=18, spacing=1.5)
code(s, Inches(0.8), Inches(2.2), Inches(11.5), Inches(2.5),
     'from typing_extensions import TypedDict, Annotated\n'
     'from langgraph.graph.message import add_messages\n'
     'from langchain_core.messages import AnyMessage\n'
     '\n'
     'class ChatState(TypedDict):\n'
     '    messages: Annotated[list[AnyMessage], add_messages]  # 메시지 누적\n'
     '    context: dict   # 공유 임시 데이터', size=15)
ml(s, Inches(0.8), Inches(5.0), Inches(11.5), Inches(2.0), [
    "핵심 포인트:",
    "  - Annotated[list, add_messages] = 메시지가 덮어쓰기가 아닌 \"누적\"됨",
    "  - add_messages는 LangGraph의 특별한 리듀서(reducer) 함수",
    "  - context: 노드 간 공유할 임시 데이터 저장소",
], size=17, spacing=1.4)


# ================================================================
# 슬라이드 8: interrupt() 사용
# ================================================================
s = new_slide()
title_line(s, "핵심 코드 2: interrupt() 로 실행 멈추기", "8")
code(s, Inches(0.8), Inches(1.3), Inches(11.5), Inches(3.5),
     'from langgraph.types import interrupt\n'
     'from langgraph.checkpoint.memory import InMemorySaver\n'
     '\n'
     'def agent_node(state: ChatState):\n'
     '    ctx = state.get("context", {})\n'
     '    food_category = ctx.get("food_category")\n'
     '\n'
     '    if not food_category:\n'
     '        # 여기서 실행이 멈춤! 사용자에게 질문을 보냄\n'
     '        food_category = interrupt("어떤 종류의 음식을 원하시나요?")\n'
     '\n'
     '    # 사용자가 답변하면 여기서부터 이어서 실행\n'
     '    messages = [system_prompt] + state["messages"]\n'
     '    response = model.invoke(messages)\n'
     '    return {"messages": [response],\n'
     '            "context": {"food_category": food_category}}', size=14)
ml(s, Inches(0.8), Inches(5.2), Inches(11.5), Inches(1.8), [
    "interrupt(질문)의 동작:",
    "  1. 그래프 실행이 즉시 멈춤 (Pause)",
    "  2. 괄호 안의 질문이 사용자에게 전달됨",
    "  3. 사용자가 답변하면 그 값이 food_category에 들어감",
], size=17, spacing=1.3)


# ================================================================
# 슬라이드 9: 실행 & 재개
# ================================================================
s = new_slide()
title_line(s, "핵심 코드 3: 실행, 상태 확인, 재개", "9")
code(s, Inches(0.8), Inches(1.3), Inches(5.7), Inches(3.0),
     '# 1. 그래프 컴파일 (체크포인터 필수!)\n'
     'memory = InMemorySaver()\n'
     'app = workflow.compile(\n'
     '    checkpointer=memory\n'
     ')\n'
     '\n'
     '# 2. 실행 (thread_id로 대화 추적)\n'
     'config = {"configurable":\n'
     '          {"thread_id": "1"}}\n'
     'result = app.invoke(msg, config=config)', size=13)
code(s, Inches(6.8), Inches(1.3), Inches(5.7), Inches(3.0),
     '# 3. 멈춘 상태 확인\n'
     'snapshot = app.get_state(config)\n'
     'if snapshot.next:  # 멈춰있다면\n'
     '    question = snapshot.tasks[0]\\\n'
     '        .interrupts[0].value\n'
     '    print(question)  # 질문 출력\n'
     '\n'
     '# 4. 답변 보내고 재개!\n'
     'from langgraph.types import Command\n'
     'result = app.invoke(\n'
     '    Command(resume="일식"), config=config\n'
     ')', size=13)
ml(s, Inches(0.8), Inches(4.7), Inches(11.5), Inches(2.3), [
    "핵심 정리:",
    "  - InMemorySaver = 멈춘 시점의 상태를 메모리에 저장",
    "  - thread_id = 대화를 구분하는 고유 ID (여러 대화 동시 관리 가능)",
    "  - get_state() = 현재 상태 확인 (멈춰있는지, 어떤 질문인지)",
    "  - Command(resume=값) = 답변을 보내면서 실행 재개",
], size=17, spacing=1.3)


# ================================================================
# 슬라이드 10: SECTION 2 구분
# ================================================================
section_slide("SECTION 2", "Interrupts in Tools",
              "도구(Tool) 내부에서 사용자에게 질문하는 패턴")


# ================================================================
# 슬라이드 11: Interrupts in Tools 이론 + 코드
# ================================================================
s = new_slide()
title_line(s, "도구(Tool) 안에서 interrupt 사용하기", "11")
ml(s, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.8), [
    "LLM이 도구를 호출할 때, 도구 내부에서 사용자 입력을 요청할 수 있습니다.",
    "예: AI가 \"사람에게 물어보기\" 도구를 사용 → 사용자에게 질문 전달",
], size=18, spacing=1.5)
code(s, Inches(0.8), Inches(2.5), Inches(5.7), Inches(2.5),
     'from langchain_core.tools import tool\n'
     '\n'
     '@tool\n'
     'def ask_human(question: str) -> str:\n'
     '    """사람에게 질문하는 도구"""\n'
     '    # 도구 실행 중 멈추고 질문!\n'
     '    response = interrupt(question)\n'
     '    return response', size=14)
code(s, Inches(6.8), Inches(2.5), Inches(5.7), Inches(2.5),
     '# 조건부 라우터: 도구 호출 여부 판단\n'
     'def should_continue(state):\n'
     '    last = state["messages"][-1]\n'
     '    if last.tool_calls:\n'
     '        return "tool_node"  # 도구 실행\n'
     '    return END  # 종료\n'
     '\n'
     'workflow.add_conditional_edges(\n'
     '    "agent", should_continue)', size=14)
ml(s, Inches(0.8), Inches(5.3), Inches(11.5), Inches(1.8), [
    "실행 흐름:",
    "  agent → LLM이 ask_human 도구 호출 → tool_node 실행 →",
    "  interrupt(질문) → 멈춤 → Command(resume=답변) → 도구 결과 반환 → agent",
], size=17, spacing=1.4)


# ================================================================
# 슬라이드 12: Interrupts in Tools 정리
# ================================================================
s = new_slide()
title_line(s, "실행 흐름 정리: Interrupts in Tools", "12")
ml(s, Inches(0.8), Inches(1.4), Inches(11.5), Inches(5.5), [
    "Section 1 vs Section 2 비교:",
    "",
    "  [Section 1] 노드에서 직접 interrupt",
    "    → 개발자가 \"여기서 멈춰!\"라고 명시적으로 지정",
    "    → 항상 같은 시점에서 멈춤",
    "",
    "  [Section 2] 도구 안에서 interrupt",
    "    → LLM이 \"사람에게 물어봐야겠다\"고 판단하면 멈춤",
    "    → AI가 스스로 언제 질문할지 결정 (더 유연함!)",
    "",
    "공통점:",
    "  - 둘 다 InMemorySaver(체크포인터)가 필수",
    "  - 둘 다 Command(resume=값)으로 재개",
    "  - 둘 다 thread_id로 대화를 추적",
], size=19, spacing=1.3)


# ================================================================
# 슬라이드 13: SECTION 3 구분
# ================================================================
section_slide("SECTION 3", "Validating Human Input",
              "사용자 입력을 검증하고 올바른 값을 받을 때까지 반복")


# ================================================================
# 슬라이드 14: 입력 검증 이론 + 코드
# ================================================================
s = new_slide()
title_line(s, "입력 검증 루프 패턴", "14")
ml(s, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.8), [
    "사용자가 잘못된 값을 입력하면? → 다시 물어보는 루프를 만듭니다.",
    "비유: 비밀번호 설정 시 \"8자 이상 입력하세요\" → 조건 충족할 때까지 반복",
], size=18, spacing=1.5)
code(s, Inches(0.8), Inches(2.5), Inches(11.5), Inches(3.0),
     'def human_node(state: ChatState):\n'
     '    """사용자 입력을 검증하는 노드"""\n'
     '    # 첫 번째 질문\n'
     '    result = interrupt("선택하세요 (A/B/C):")\n'
     '\n'
     '    # 올바른 값이 나올 때까지 반복 질문!\n'
     '    while result not in ["A", "B", "C"]:\n'
     '        result = interrupt(\n'
     '            f"\\"{result}\\\"는 잘못된 입력입니다! A, B, C 중 선택하세요:"\n'
     '        )\n'
     '\n'
     '    return {"messages": [HumanMessage(content=result)]}', size=15)
ml(s, Inches(0.8), Inches(5.8), Inches(11.5), Inches(1.2), [
    "핵심: interrupt()를 while 루프 안에 넣으면 검증 루프 완성!",
    "매번 멈추고 → 사용자 입력 → 검증 → 실패시 다시 멈춤",
], size=17, spacing=1.4)


# ================================================================
# 슬라이드 15: 입력 검증 정리
# ================================================================
s = new_slide()
title_line(s, "정리: Interrupt 3가지 활용법", "15")
ml(s, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5), [
    "1. 노드에서 직접 interrupt (Section 1)",
    "   → 개발자가 정한 시점에 항상 멈춤",
    "   → 사용 예: 주문 확인, 필수 정보 수집",
    "",
    "2. 도구 안에서 interrupt (Section 2)",
    "   → AI가 판단해서 필요할 때만 멈춤",
    "   → 사용 예: 대화 중 추가 정보가 필요할 때",
    "",
    "3. 검증 루프와 interrupt (Section 3)",
    "   → 올바른 입력이 올 때까지 반복",
    "   → 사용 예: 선택지 검증, 형식 검사",
    "",
    "공통 필수 요소: InMemorySaver + thread_id + Command(resume=값)",
], size=19, spacing=1.3)


# ================================================================
# 슬라이드 16: SECTION 4 구분
# ================================================================
section_slide("SECTION 4", "Command(goto=)",
              "런타임에 다음 실행 노드를 동적으로 결정하는 패턴")


# ================================================================
# 슬라이드 17: Command goto 이론
# ================================================================
s = new_slide()
title_line(s, "Command(goto=) vs add_conditional_edges()", "17")
ml(s, Inches(0.8), Inches(1.4), Inches(11.5), Inches(1.5), [
    "\"다음에 어디로 갈지\"를 결정하는 두 가지 방법 비교:",
], size=20, spacing=1.5)
ml(s, Inches(0.8), Inches(2.5), Inches(11.5), Inches(4.5), [
    "┌─────────────────────┬────────────────────────┬──────────────────────────┐",
    "│                     │ add_conditional_edges  │ Command(goto=)           │",
    "├─────────────────────┼────────────────────────┼──────────────────────────┤",
    "│ 결정 시점           │ 그래프 빌드 시         │ 런타임 (실행 중)         │",
    "│ 유연성              │ 미리 정한 경우만 가능  │ 완전 자유 (동적)         │",
    "│ 사용 위치           │ 그래프 구성 코드       │ 노드 함수 내부           │",
    "│ 결과와 동시 업데이트│ 불가                   │ Command(update=, goto=)  │",
    "└─────────────────────┴────────────────────────┴──────────────────────────┘",
], size=14, spacing=1.0, font="Consolas")
ml(s, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.5), [
    "비유:",
    "  add_conditional_edges = 미리 정해둔 교차로 표지판 (좌/우만 가능)",
    "  Command(goto=) = 내비게이션 (실시간으로 최적 경로 안내)",
], size=18, spacing=1.4)


# ================================================================
# 슬라이드 18: Command goto 핵심 코드
# ================================================================
s = new_slide()
title_line(s, "핵심 코드: 레스토랑 예약 흐름", "18")
tb(s, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.3),
   "추천 → 승인 → 예약 (거절 시 다시 추천으로 돌아감)", size=18)
code(s, Inches(0.8), Inches(1.8), Inches(11.5), Inches(4.5),
     'def human_approval_node(state: ChatState):\n'
     '    menu = state["selected_menu"]\n'
     '    # 사용자에게 승인 요청\n'
     '    approved = interrupt(f"\\"{menu}\\\" 예약할까요? (yes/no)")\n'
     '\n'
     '    if approved == "yes":\n'
     '        return Command(\n'
     '            update={"messages": ["승인됨"]},\n'
     '            goto="booking"      # 예약 노드로 이동!\n'
     '        )\n'
     '    else:\n'
     '        return Command(\n'
     '            update={"messages": ["거절됨"]},\n'
     '            goto="recommender"  # 다시 추천 노드로 돌아감!\n'
     '        )', size=14)
ml(s, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.7), [
    "Command(update=, goto=) = State 업데이트 + 다음 노드 이동을 한 번에!",
], size=17, spacing=1.4)


# ================================================================
# 슬라이드 19: Command goto 그래프 & 정리
# ================================================================
s = new_slide()
title_line(s, "그래프 구성 & 정리: Command(goto=)", "19")
code(s, Inches(0.8), Inches(1.3), Inches(5.7), Inches(2.8),
     'workflow = StateGraph(ChatState)\n'
     'workflow.add_node("recommender",\n'
     '                  recommender_node)\n'
     'workflow.add_node("approval",\n'
     '                  human_approval_node)\n'
     'workflow.add_node("booking",\n'
     '                  booking_node)\n'
     '\n'
     'workflow.add_edge(START, "recommender")\n'
     'workflow.add_edge("recommender", "approval")\n'
     '# booking으로의 엣지는 없음!\n'
     '# → Command(goto)로 런타임에 결정', size=13)
ml(s, Inches(6.8), Inches(1.3), Inches(5.7), Inches(2.8), [
    "실행 흐름:",
    "",
    "START → recommender → approval",
    "  ├─ yes → Command(goto='booking')",
    "  │        → booking → END",
    "  └─ no  → Command(goto='recommender')",
    "           → recommender → approval ...",
    "",
    "엣지 없이도 노드 이동 가능!",
], size=16, spacing=1.2)
ml(s, Inches(0.8), Inches(4.5), Inches(11.5), Inches(2.5), [
    "Command(goto) 핵심 정리:",
    "  - 노드 함수 내부에서 return Command(goto='노드명')으로 이동",
    "  - update= 로 State 업데이트도 동시에 가능",
    "  - 미리 엣지를 정의하지 않아도 되므로 매우 유연",
    "  - interrupt()와 결합하면 \"사용자 결정 → 동적 분기\" 패턴 완성!",
], size=17, spacing=1.3)


# ================================================================
# 슬라이드 20: SECTION 5 구분
# ================================================================
section_slide("SECTION 5", "Subgraph (State 공유)",
              "그래프 안에 그래프를 넣어 모듈화하는 패턴")


# ================================================================
# 슬라이드 21: Subgraph 이론
# ================================================================
s = new_slide()
title_line(s, "Subgraph란?", "21")
ml(s, Inches(0.8), Inches(1.4), Inches(11.5), Inches(2.0), [
    "하나의 그래프 안에 다른 그래프를 \"노드처럼\" 넣는 것",
    "",
    "비유: 레고 블록",
    "  작은 레고(서브그래프)를 미리 조립 → 큰 레고(부모 그래프)에 끼워 넣기",
], size=20, spacing=1.4)
ml(s, Inches(0.8), Inches(3.5), Inches(11.5), Inches(3.5), [
    "두 가지 연결 방식:",
    "",
    "  1. State 공유형 (이번 섹션)",
    "     부모와 자식이 같은 State 사용",
    "     workflow.add_node('kitchen', kitchen_subgraph)  # 서브그래프를 노드로!",
    "",
    "  2. State 분리형 (다음 섹션)",
    "     부모와 자식이 다른 State 사용",
    "     부모 노드 함수에서 직접 자식 그래프를 invoke",
    "",
    "예제: 레스토랑(부모) + 주방(자식 서브그래프)",
], size=18, spacing=1.2)


# ================================================================
# 슬라이드 22: Subgraph State 공유 핵심 코드
# ================================================================
s = new_slide()
title_line(s, "핵심 코드: 레스토랑 + 주방 서브그래프", "22")
tb(s, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.3),
   "자식: 주방 서브그래프", size=16, bold=True)
code(s, Inches(0.8), Inches(1.7), Inches(5.7), Inches(3.0),
     '# 주방 서브그래프 (같은 State 사용)\n'
     'def check_ingredients(state):\n'
     '    print(f"재고 확인: {state[\'menu\']}")\n'
     '    return {"status": "ok"}\n'
     '\n'
     'def assign_chef(state):\n'
     '    if "오마카세" in state["menu"]:\n'
     '        return {"chef_name": "Master Jiro",\n'
     '                "cooking_time": "40분"}\n'
     '    return {"chef_name": "Chef Kim",\n'
     '            "cooking_time": "20분"}\n'
     '\n'
     'kitchen = StateGraph(ReservationState)\n'
     'kitchen.add_node("check", check_ingredients)\n'
     'kitchen.add_node("assign", assign_chef)\n'
     '# ... edges, compile', size=11)

tb(s, Inches(6.8), Inches(1.3), Inches(5.5), Inches(0.3),
   "부모: 레스토랑 그래프", size=16, bold=True)
code(s, Inches(6.8), Inches(1.7), Inches(5.7), Inches(3.0),
     '# 부모 그래프\n'
     'parent = StateGraph(ReservationState)\n'
     'parent.add_node("menu_recommend",\n'
     '                menu_recommender)\n'
     '\n'
     '# 서브그래프를 노드로 등록!\n'
     'parent.add_node("kitchen",\n'
     '                kitchen_subgraph)\n'
     '\n'
     'parent.add_edge(START, "menu_recommend")\n'
     'parent.add_edge("menu_recommend",\n'
     '                "kitchen")\n'
     'parent.add_edge("kitchen", END)\n'
     '\n'
     'app = parent.compile()', size=11)

ml(s, Inches(0.8), Inches(5.0), Inches(11.5), Inches(2.0), [
    "핵심: add_node('이름', 컴파일된_서브그래프) — 서브그래프를 노드처럼 사용!",
    "State 공유형은 부모와 자식이 같은 ReservationState를 사용하므로 데이터가 자동 전달됨",
], size=17, spacing=1.4)


# ================================================================
# 슬라이드 23: Subgraph State 공유 정리
# ================================================================
s = new_slide()
title_line(s, "실행 흐름 & 정리: Subgraph (State 공유)", "23")
ml(s, Inches(0.8), Inches(1.4), Inches(11.5), Inches(2.5), [
    "실행 흐름:",
    "",
    "  START → menu_recommend → [kitchen 서브그래프]",
    "                                ├→ check_ingredients",
    "                                └→ assign_chef",
    "                            → END",
], size=18, spacing=1.3, font="Consolas")
ml(s, Inches(0.8), Inches(4.0), Inches(11.5), Inches(3.0), [
    "State 공유형 장점:",
    "  - 구현이 간단 (add_node 한 줄이면 끝)",
    "  - State가 자동으로 부모 ↔ 자식 간 전달",
    "  - 서브그래프 내부의 노드 실행 결과가 바로 부모 State에 반영",
    "",
    "State 공유형 단점:",
    "  - 부모와 자식이 같은 State를 써야 함 (결합도가 높음)",
    "  - 외부 서비스 등 독립적인 시스템과는 맞지 않을 수 있음",
], size=18, spacing=1.3)


# ================================================================
# 슬라이드 24: SECTION 6 구분
# ================================================================
section_slide("SECTION 6", "Subgraph (State 분리)",
              "부모와 자식이 서로 다른 State를 사용하는 패턴")


# ================================================================
# 슬라이드 25: State 분리형 이론
# ================================================================
s = new_slide()
title_line(s, "왜 State를 분리해야 하는가?", "25")
ml(s, Inches(0.8), Inches(1.4), Inches(11.5), Inches(5.5), [
    "외부 시스템은 우리 State 구조를 모릅니다!",
    "",
    "비유: 레스토랑에서 배달 업체에 주문할 때",
    "  레스토랑: 메뉴명, 고객 주소, 테이블 번호, 결제 정보...",
    "  배달 업체: 주소, 물건 정보, 배달 상태만 필요",
    "  → 레스토랑 정보를 배달 업체 양식에 맞춰 \"변환\"해서 전달!",
    "",
    "State 분리가 필요한 경우:",
    "  - 외부 API, 결제 시스템, 배송 시스템과 연동",
    "  - 독립적으로 개발된 서브시스템 통합",
    "  - 보안상 일부 데이터만 전달해야 할 때",
    "",
    "핵심 아이디어: \"브릿지 노드\"가 부모 State ↔ 자식 State 변환을 담당",
], size=18, spacing=1.3)


# ================================================================
# 슬라이드 26: State 분리형 핵심 코드
# ================================================================
s = new_slide()
title_line(s, "핵심 코드: 브릿지 노드 패턴", "26")
code(s, Inches(0.8), Inches(1.3), Inches(5.7), Inches(2.0),
     '# 부모 State (레스토랑)\n'
     'class RestaurantState(TypedDict):\n'
     '    menu: str\n'
     '    customer_address: str\n'
     '    final_status: str', size=14)
code(s, Inches(6.8), Inches(1.3), Inches(5.7), Inches(2.0),
     '# 자식 State (배달 업체)\n'
     'class DeliveryState(TypedDict):\n'
     '    address: str\n'
     '    package_info: str\n'
     '    delivery_status: str', size=14)
code(s, Inches(0.8), Inches(3.5), Inches(11.5), Inches(3.3),
     '# 브릿지 노드: 부모 State → 자식 State → 부모 State\n'
     'def call_delivery_service(state: RestaurantState):\n'
     '    # Step 1: 부모 → 자식 변환 (필요한 정보만 전달)\n'
     '    delivery_input = {\n'
     '        "address": state["customer_address"],  # 부모의 customer_address → 자식의 address\n'
     '        "package_info": state["menu"],          # 부모의 menu → 자식의 package_info\n'
     '        "delivery_status": "pending"\n'
     '    }\n'
     '\n'
     '    # Step 2: 자식 그래프 실행\n'
     '    result = delivery_graph.invoke(delivery_input)\n'
     '\n'
     '    # Step 3: 자식 → 부모 변환 (결과만 가져옴)\n'
     '    return {"final_status": result["delivery_status"]}', size=13)


# ================================================================
# 슬라이드 27: State 분리형 정리
# ================================================================
s = new_slide()
title_line(s, "정리: Subgraph 두 가지 방식 비교", "27")
ml(s, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5), [
    "┌──────────────────┬──────────────────────────┬──────────────────────────┐",
    "│                  │ State 공유형             │ State 분리형             │",
    "├──────────────────┼──────────────────────────┼──────────────────────────┤",
    "│ State 구조       │ 부모 = 자식 (동일)       │ 부모 =/= 자식 (다름)     │",
    "│ 연결 방법        │ add_node(name, subgraph) │ 노드 함수에서 invoke     │",
    "│ 데이터 전달      │ 자동                     │ 수동 (브릿지 노드)       │",
    "│ 장점             │ 간단, 코드 적음          │ 독립성, 보안, 유연성     │",
    "│ 사용 사례        │ 내부 모듈 분리           │ 외부 시스템 연동         │",
    "└──────────────────┴──────────────────────────┴──────────────────────────┘",
], size=14, spacing=1.0, font="Consolas")
ml(s, Inches(0.8), Inches(5.0), Inches(11.5), Inches(2.0), [
    "선택 기준:",
    "  - 같은 팀이 만든 모듈? → State 공유형 (간단)",
    "  - 외부 시스템이나 독립 서비스? → State 분리형 (안전)",
], size=19, spacing=1.5)


# ================================================================
# 슬라이드 28: SECTION 7 구분
# ================================================================
section_slide("SECTION 7", "Multi-Agent: Supervisor",
              "중앙 관리자가 워커 에이전트를 지휘하는 패턴")


# ================================================================
# 슬라이드 29: Supervisor 이론
# ================================================================
s = new_slide()
title_line(s, "Supervisor 패턴이란?", "29")
ml(s, Inches(0.8), Inches(1.4), Inches(11.5), Inches(2.0), [
    "하나의 \"관리자(Supervisor)\"가 여러 \"작업자(Worker)\"를 지휘하는 구조",
    "",
    "비유: 편집장과 기자들",
    "  편집장 → \"기자A, 조사해와\" → 기자A 결과 확인 → \"기자B, 글 써\" → 기자B 결과 확인 → 완료!",
], size=20, spacing=1.4)
ml(s, Inches(0.8), Inches(3.5), Inches(11.5), Inches(3.5), [
    "아키텍처:",
    "",
    "  사용자 요청",
    "      ↓",
    "  [Supervisor] ──→ [Researcher] → 조사 → [Supervisor]",
    "               ──→ [Writer]     → 작성 → [Supervisor]",
    "               ──→ FINISH (완료!)",
    "",
    "핵심: 모든 Worker는 작업 후 반드시 Supervisor에게 돌아감!",
    "Supervisor가 다음에 누가 일할지 매번 결정 (중앙 집중식)",
], size=18, spacing=1.2)


# ================================================================
# 슬라이드 30: Supervisor 핵심 코드 1
# ================================================================
s = new_slide()
title_line(s, "핵심 코드 1: Supervisor 노드", "30")
code(s, Inches(0.8), Inches(1.3), Inches(11.5), Inches(5.0),
     'from pydantic import BaseModel, Field\n'
     'from typing import Literal\n'
     '\n'
     '# Supervisor의 결정을 구조화\n'
     'class RouteResponse(BaseModel):\n'
     '    next_step: Literal["researcher", "writer", "FINISH"]\n'
     '    reason: str\n'
     '\n'
     'supervisor_llm = model.with_structured_output(RouteResponse)\n'
     '\n'
     'def supervisor_node(state: TeamState):\n'
     '    system_prompt = """\n'
     '    당신은 블로그 작성 팀 관리자입니다.\n'
     '    1. researcher에게 조사를 시킵니다\n'
     '    2. writer에게 글 작성을 시킵니다\n'
     '    3. 완료되면 FINISH를 선택합니다\n'
     '    """\n'
     '    messages = [SystemMessage(content=system_prompt)] + state["messages"]\n'
     '    decision = supervisor_llm.invoke(messages)\n'
     '\n'
     '    if decision.next_step == "FINISH":\n'
     '        return Command(goto=END)           # 종료!\n'
     '    else:\n'
     '        return Command(goto=decision.next_step)  # 해당 Worker로 이동', size=13)


# ================================================================
# 슬라이드 31: Supervisor 핵심 코드 2
# ================================================================
s = new_slide()
title_line(s, "핵심 코드 2: Worker 노드", "31")
code(s, Inches(0.8), Inches(1.3), Inches(5.7), Inches(3.5),
     '# Researcher 노드\n'
     'def researcher_node(state: TeamState):\n'
     '    system_msg = SystemMessage(\n'
     '        content="IT 트렌드 전문 리서처"\n'
     '    )\n'
     '    messages = [system_msg] \\\n'
     '              + state["messages"]\n'
     '    response = model.invoke(messages)\n'
     '\n'
     '    return Command(\n'
     '        update={"messages":\n'
     '          [HumanMessage(\n'
     '            content=response.content,\n'
     '            name="Researcher")]},\n'
     '        goto="supervisor"  # 관리자에게 복귀!\n'
     '    )', size=12)
code(s, Inches(6.8), Inches(1.3), Inches(5.7), Inches(3.5),
     '# Writer 노드\n'
     'def writer_node(state: TeamState):\n'
     '    system_msg = SystemMessage(\n'
     '        content="테크 블로그 전문 작가"\n'
     '    )\n'
     '    messages = [system_msg] \\\n'
     '              + state["messages"]\n'
     '    response = model.invoke(messages)\n'
     '\n'
     '    return Command(\n'
     '        update={"messages":\n'
     '          [HumanMessage(\n'
     '            content=response.content,\n'
     '            name="Writer")]},\n'
     '        goto="supervisor"  # 관리자에게 복귀!\n'
     '    )', size=12)
ml(s, Inches(0.8), Inches(5.1), Inches(11.5), Inches(2.0), [
    "핵심 패턴:",
    "  - Worker는 작업 후 항상 goto='supervisor'로 복귀",
    "  - name='Researcher' / 'Writer' — Supervisor가 누가 한 작업인지 구분",
    "  - HumanMessage로 보내는 이유: 다음 LLM 호출에서 \"이전 대화\"로 인식됨",
], size=17, spacing=1.3)


# ================================================================
# 슬라이드 32: Supervisor 실행 & 정리
# ================================================================
s = new_slide()
title_line(s, "실행 & 정리: Supervisor 패턴", "32")
code(s, Inches(0.8), Inches(1.3), Inches(11.5), Inches(2.3),
     '# 그래프 구성\n'
     'workflow = StateGraph(TeamState)\n'
     'workflow.add_node("supervisor", supervisor_node)\n'
     'workflow.add_node("researcher", researcher_node)\n'
     'workflow.add_node("writer", writer_node)\n'
     'workflow.add_edge(START, "supervisor")  # 항상 Supervisor부터 시작!\n'
     '\n'
     '# 실행\n'
     'result = app.invoke({"messages": [HumanMessage("AI 트렌드 블로그 작성해줘")]})', size=14)
ml(s, Inches(0.8), Inches(4.0), Inches(11.5), Inches(3.0), [
    "실행 흐름 예시:",
    "  1. Supervisor: \"먼저 조사가 필요\" → goto='researcher'",
    "  2. Researcher: 조사 완료 → goto='supervisor'",
    "  3. Supervisor: \"이제 글 작성\" → goto='writer'",
    "  4. Writer: 작성 완료 → goto='supervisor'",
    "  5. Supervisor: \"다 됐다\" → goto=END (FINISH)",
], size=18, spacing=1.3)


# ================================================================
# 슬라이드 33: SECTION 8 구분
# ================================================================
section_slide("SECTION 8", "Multi-Agent: Handoff",
              "에이전트끼리 직접 업무를 넘기는 분산형 패턴")


# ================================================================
# 슬라이드 34: Handoff 이론
# ================================================================
s = new_slide()
title_line(s, "Handoff 패턴이란?", "34")
ml(s, Inches(0.8), Inches(1.4), Inches(11.5), Inches(2.0), [
    "중앙 관리자 없이 에이전트끼리 직접 업무를 넘기는 분산형 구조",
    "",
    "비유: 고객센터 전화 연결",
    "  안내 데스크 → \"기술 문제군요, 기술팀으로 연결해드릴게요\" → 기술팀 직접 응대",
], size=20, spacing=1.4)
ml(s, Inches(0.8), Inches(3.5), Inches(11.5), Inches(3.5), [
    "아키텍처:",
    "",
    "  [Receptionist] ──(기술 문의)──→ [Tech Support]",
    "                 ←──(이관)──────────────│",
    "                                        │",
    "                 ──(결제 문의)──→ [Billing]",
    "                 ←──(이관)──────────│",
    "",
    "핵심: 각 에이전트가 자기 전용 \"핸드오프 도구\"를 가지고 있음!",
    "도구를 호출하면 → Command(goto='다른에이전트')로 업무 이관",
], size=18, spacing=1.2)


# ================================================================
# 슬라이드 35: Handoff Tools 핵심 코드
# ================================================================
s = new_slide()
title_line(s, "핵심 코드 1: Handoff 도구 (Transfer Tools)", "35")
code(s, Inches(0.8), Inches(1.3), Inches(11.5), Inches(5.0),
     'from langchain_core.tools import tool\n'
     'from langchain_core.messages import ToolMessage\n'
     '\n'
     '@tool\n'
     'def transfer_to_tech_support(runtime):\n'
     '    """기술 문제를 기술팀으로 연결하는 도구"""\n'
     '    tool_id = runtime.tool_call_id\n'
     '    return Command(\n'
     '        update={"messages": [\n'
     '            ToolMessage(content="기술팀으로 연결됨",\n'
     '                        tool_call_id=tool_id)\n'
     '        ]},\n'
     '        goto="tech_support"       # 기술팀 에이전트로 이동!\n'
     '    )\n'
     '\n'
     '@tool\n'
     'def transfer_to_billing(runtime):\n'
     '    """결제/환불 문제를 환불팀으로 연결하는 도구"""\n'
     '    tool_id = runtime.tool_call_id\n'
     '    return Command(\n'
     '        update={"messages": [\n'
     '            ToolMessage(content="환불팀으로 연결됨",\n'
     '                        tool_call_id=tool_id)\n'
     '        ]},\n'
     '        goto="billing"            # 환불팀 에이전트로 이동!\n'
     '    )', size=13)


# ================================================================
# 슬라이드 36: Agent 노드 핵심 코드
# ================================================================
s = new_slide()
title_line(s, "핵심 코드 2: 에이전트 노드", "36")
code(s, Inches(0.8), Inches(1.3), Inches(5.7), Inches(3.0),
     '# 안내 데스크 (시작점)\n'
     'def receptionist_node(state):\n'
     '    system = "고객 분류 후\\n"\\\n'
     '        "기술문의→transfer_to_tech\\n"\\\n'
     '        "결제문의→transfer_to_billing"\n'
     '    tools = [transfer_to_tech_support,\n'
     '             transfer_to_billing]\n'
     '    llm = model.bind_tools(tools)\n'
     '    msgs = [SystemMessage(content=system)]\\\n'
     '           + state["messages"]\n'
     '    return {"messages": [llm.invoke(msgs)]}', size=12)
code(s, Inches(6.8), Inches(1.3), Inches(5.7), Inches(3.0),
     '# 기술팀 에이전트\n'
     'def tech_node(state):\n'
     '    system = "기술 문제 해결\\n"\\\n'
     '        "환불→transfer_to_billing\\n"\\\n'
     '        "해결→transfer_to_reception"\n'
     '    tools = [transfer_to_billing,\n'
     '             transfer_to_reception]\n'
     '    llm = model.bind_tools(tools)\n'
     '    msgs = [SystemMessage(content=system)]\\\n'
     '           + state["messages"]\n'
     '    return {"messages": [llm.invoke(msgs)]}', size=12)
ml(s, Inches(0.8), Inches(4.6), Inches(11.5), Inches(2.5), [
    "핵심 패턴:",
    "  - 각 에이전트는 자기가 연결할 수 있는 도구만 가짐",
    "  - bind_tools(tools) — LLM에 사용 가능한 도구 목록을 알려줌",
    "  - LLM이 상황을 판단해서 적절한 transfer 도구를 호출",
    "  - Supervisor 없이도 자율적으로 업무 이관 가능!",
], size=17, spacing=1.3)


# ================================================================
# 슬라이드 37: 그래프 조립 핵심 코드
# ================================================================
s = new_slide()
title_line(s, "핵심 코드 3: 그래프 조립 (ToolNode)", "37")
code(s, Inches(0.8), Inches(1.3), Inches(11.5), Inches(4.5),
     'from langgraph.prebuilt import ToolNode\n'
     '\n'
     '# 모든 도구를 ToolNode에 등록\n'
     'all_tools = [transfer_to_tech_support,\n'
     '             transfer_to_billing,\n'
     '             transfer_to_reception]\n'
     'tool_node = ToolNode(all_tools)\n'
     '\n'
     '# 그래프 구성\n'
     'workflow = StateGraph(AgentState)\n'
     'workflow.add_node("receptionist", receptionist_node)\n'
     'workflow.add_node("tech_support", tech_node)\n'
     'workflow.add_node("billing", billing_node)\n'
     'workflow.add_node("tools", tool_node)    # 도구 실행 노드\n'
     '\n'
     'workflow.add_edge(START, "receptionist")  # 항상 안내 데스크부터 시작\n'
     '\n'
     '# 각 에이전트 → tools (도구 실행)\n'
     'workflow.add_edge("receptionist", "tools")\n'
     'workflow.add_edge("tech_support", "tools")\n'
     'workflow.add_edge("billing", "tools")\n'
     '\n'
     'app = workflow.compile()', size=13)
ml(s, Inches(0.8), Inches(6.0), Inches(11.5), Inches(1.0), [
    "ToolNode = 도구를 실행해주는 내장 노드 (도구 안의 Command(goto)가 다음 경로를 결정!)",
], size=17, spacing=1.4)


# ================================================================
# 슬라이드 38: Handoff 실행 & 정리
# ================================================================
s = new_slide()
title_line(s, "실행 & 정리: Handoff 패턴", "38")
code(s, Inches(0.8), Inches(1.3), Inches(11.5), Inches(1.5),
     '# 실행 예시\n'
     'result = app.invoke({\n'
     '    "messages": [HumanMessage("노트북 충전이 안 돼요. 어제 산건데 환불도 가능한가요?")]\n'
     '})', size=14)
ml(s, Inches(0.8), Inches(3.2), Inches(11.5), Inches(4.0), [
    "예상 흐름:",
    "  1. receptionist: \"기술 문제군요\" → transfer_to_tech_support 도구 호출",
    "  2. tools: 도구 실행 → Command(goto='tech_support')",
    "  3. tech_support: \"충전 문제 안내... 환불은 환불팀으로\" → transfer_to_billing",
    "  4. tools: 도구 실행 → Command(goto='billing')",
    "  5. billing: \"환불 절차 안내\" → transfer_to_reception",
    "  6. tools: 도구 실행 → Command(goto='receptionist')",
    "  7. receptionist: \"더 도와드릴 것이 있나요?\"",
], size=17, spacing=1.3)


# ================================================================
# 슬라이드 39: Supervisor vs Handoff 비교표
# ================================================================
s = new_slide()
title_line(s, "Supervisor vs Handoff 비교", "39")
ml(s, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5), [
    "┌──────────────────┬──────────────────────────┬──────────────────────────┐",
    "│                  │ Supervisor (중앙 집중형) │ Handoff (분산형)         │",
    "├──────────────────┼──────────────────────────┼──────────────────────────┤",
    "│ 구조             │ 관리자가 Worker를 지휘   │ Agent끼리 직접 이관      │",
    "│ 비유             │ 편집장 + 기자들          │ 고객센터 전화 연결       │",
    "│ 라우팅 방식      │ Supervisor LLM 판단      │ 각 Agent의 도구 호출     │",
    "│ Worker 복귀      │ 항상 Supervisor로 복귀   │ 다음 Agent로 직접 이동   │",
    "│ 전체 상황 파악   │ Supervisor가 전체 파악   │ 각자 자기 영역만 파악    │",
    "│ 유연성           │ 중간 (관리자 의존)       │ 높음 (자율적)            │",
    "│ 사용 사례        │ 팀 프로젝트, 순서 중요   │ 고객센터, 독립 부서 연동 │",
    "│ 핵심 API         │ with_structured_output   │ bind_tools + ToolNode    │",
    "│                  │ + Command(goto)          │ + Command(goto)          │",
    "└──────────────────┴──────────────────────────┴──────────────────────────┘",
], size=14, spacing=1.0, font="Consolas")


# ================================================================
# 슬라이드 40: 8개 패턴 총정리
# ================================================================
s = new_slide()
title_line(s, "8개 패턴 한눈에 보기", "40")
ml(s, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5), [
    "[기초 - Interrupt]",
    "  1. Dynamic Interrupt      interrupt() + Command(resume=)",
    "  2. Interrupts in Tools    @tool 안에서 interrupt()",
    "  3. Validating Input       while 루프 + interrupt()",
    "",
    "[중급 - Dynamic Routing]",
    "  4. Command(goto=)         런타임에 다음 노드를 동적 결정",
    "",
    "[고급 - Subgraph]",
    "  5. State 공유형           add_node(name, subgraph)",
    "  6. State 분리형           노드 함수에서 subgraph.invoke()",
    "",
    "[심화 - Multi-Agent]",
    "  7. Supervisor             중앙 관리자 + Worker (with_structured_output)",
    "  8. Handoff                Agent 직접 이관 (bind_tools + ToolNode)",
], size=18, spacing=1.1)


# ================================================================
# 슬라이드 41: 핵심 API 총정리
# ================================================================
s = new_slide()
title_line(s, "핵심 API 총정리", "41")
ml(s, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5), [
    "interrupt(질문)",
    "  그래프 실행을 멈추고 사용자에게 질문 전달",
    "",
    "Command(resume=값)",
    "  멈춘 그래프에 답변을 보내고 실행 재개",
    "",
    "Command(goto='노드명')",
    "  런타임에 다음 실행 노드를 동적으로 결정",
    "",
    "Command(update={}, goto='노드명')",
    "  State 업데이트 + 노드 이동을 동시에 수행",
    "",
    "InMemorySaver / checkpointer",
    "  interrupt 사용 시 필수! 멈춘 시점의 상태를 저장",
    "",
    "model.with_structured_output(Schema)",
    "  LLM 출력을 Pydantic 모델로 강제 (Supervisor 라우팅에 사용)",
    "",
    "model.bind_tools(tools) / ToolNode(tools)",
    "  LLM에 도구 연결 / 도구 실행 노드 (Handoff에 사용)",
], size=17, spacing=1.05)


# ================================================================
# 슬라이드 42: Q&A
# ================================================================
s = new_slide()
tb(s, Inches(1), Inches(2.0), Inches(11), Inches(1.2),
   "Q & A", size=60, bold=True, align=PP_ALIGN.CENTER, font="Consolas")
tb(s, Inches(1), Inches(3.5), Inches(11), Inches(0.8),
   "질문이 있으신가요?", size=28, align=PP_ALIGN.CENTER)
ml(s, Inches(1), Inches(5.0), Inches(11), Inches(1.5), [
    "오늘 배운 8가지 패턴을 직접 실습해보세요!",
    "코드를 바꿔가며 실행하면 이해가 더 깊어집니다.",
], size=20, spacing=1.5, align=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────
# 저장
# ──────────────────────────────────────────────
output_path = r"c:\DL_Practice\[5강]_Multi_Agent_강의슬라이드.pptx"
prs.save(output_path)
print(f"PPT 생성 완료: {output_path}")
print(f"총 슬라이드 수: {len(prs.slides)}")
