# -*- coding: utf-8 -*-
"""2강 랭그래프 기초 - 강의 슬라이드 생성 스크립트"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ============================================================
# 헬퍼 함수
# ============================================================

BLACK = RGBColor(0, 0, 0)
DARK_GRAY = RGBColor(80, 80, 80)
CODE_BG = RGBColor(240, 240, 240)
WHITE = RGBColor(255, 255, 255)
ACCENT = RGBColor(0, 102, 204)


def add_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    return slide


def add_title_text(slide, text, top=Inches(0.4), left=Inches(0.6), width=Inches(8.8), font_size=28, bold=True, color=BLACK):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "맑은 고딕"
    return txBox


def add_subtitle_text(slide, text, top=Inches(1.1), left=Inches(0.6), width=Inches(8.8), font_size=16, color=DARK_GRAY):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.name = "맑은 고딕"
    return txBox


def add_body_text(slide, text, top=Inches(1.7), left=Inches(0.6), width=Inches(8.8), height=Inches(5.0), font_size=14, color=BLACK):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    lines = text.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "맑은 고딕"
        p.space_after = Pt(4)
    return txBox


def add_code_block(slide, code, top=Inches(3.0), left=Inches(0.6), width=Inches(8.8), height=None, font_size=11):
    lines = code.strip().split('\n')
    if height is None:
        height = Inches(0.3 + len(lines) * 0.22)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CODE_BG
    shape.line.fill.background()
    txBox = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.1), width - Inches(0.3), height - Inches(0.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = BLACK
        p.font.name = "Consolas"
        p.space_after = Pt(1)
        p.space_before = Pt(1)
    return txBox


def add_table_slide(slide, headers, rows, top=Inches(2.0), left=Inches(0.6), width=Inches(8.8)):
    row_count = len(rows) + 1
    col_count = len(headers)
    height = Inches(0.4 * row_count)
    table_shape = slide.shapes.add_table(row_count, col_count, left, top, width, height)
    table = table_shape.table
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.name = "맑은 고딕"
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(60, 60, 60)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.name = "맑은 고딕"
                p.font.color.rgb = BLACK
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE
    return table_shape


def add_flow_diagram(slide, text, top=Inches(3.5), left=Inches(0.6), width=Inches(8.8), height=Inches(1.5)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(245, 248, 255)
    shape.line.color.rgb = RGBColor(180, 200, 230)
    shape.line.width = Pt(1)
    txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), height - Inches(0.3))
    tf = txBox.text_frame
    tf.word_wrap = True
    lines = text.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(13)
        p.font.name = "Consolas"
        p.font.color.rgb = RGBColor(30, 30, 80)
        p.alignment = PP_ALIGN.CENTER
    return txBox


# ============================================================
# 프레젠테이션 생성
# ============================================================
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# ============================================================
# 슬라이드 1: 표지
# ============================================================
slide = add_slide(prs)
add_title_text(slide, "2강. LangGraph 기초", top=Inches(2.5), font_size=36, color=BLACK)
add_subtitle_text(slide, "State, Node, Edge로 AI 에이전트 구축하기", top=Inches(3.3), font_size=20, color=DARK_GRAY)
add_body_text(slide, "LangGraph 기반 AI 에이전트 프레임워크 입문", top=Inches(4.2), font_size=14, color=DARK_GRAY)

# ============================================================
# 슬라이드 2: 학습 목표
# ============================================================
slide = add_slide(prs)
add_title_text(slide, "학습 목표")
add_body_text(slide, """이번 강의를 마치면 다음을 할 수 있습니다:

1. LangGraph의 핵심 개념 (State, Node, Edge)을 설명할 수 있다
2. 도구(Tool)를 활용하는 사칙연산 에이전트를 직접 구현할 수 있다
3. 조건부 분기(Router 패턴)를 사용한 이메일 처리 에이전트를 만들 수 있다
4. 규칙 기반 분류와 LLM 기반 분류의 차이를 이해하고 적용할 수 있다

실습 내용:
  Part 1: 사칙연산 에이전트 (기초)
  Part 2: 이메일 처리 에이전트 - 규칙 기반 (Router 패턴)
  Part 3: 이메일 처리 에이전트 - LLM 기반 (개선된 분류)""", font_size=15)

# ============================================================
# 슬라이드 3: LangGraph란?
# ============================================================
slide = add_slide(prs)
add_title_text(slide, "LangGraph란?")
add_body_text(slide, """LangChain 생태계의 그래프 기반 에이전트 프레임워크

  LangChain: LLM 기반 애플리케이션을 위한 프레임워크
  LangGraph: LangChain 위에서 복잡한 워크플로우를 구현하는 도구

핵심 아이디어:
  "AI 에이전트의 실행 흐름을 그래프(Graph)로 표현한다"

  노드(Node) = 작업을 수행하는 함수
  엣지(Edge) = 노드 간의 연결선 (실행 순서)
  상태(State) = 노드들이 공유하는 데이터

왜 그래프인가?
  - 복잡한 분기/반복 로직을 시각적으로 표현 가능
  - 각 단계를 독립적으로 테스트 가능
  - 실행 흐름을 쉽게 수정/확장 가능""", font_size=14)

# ============================================================
# 슬라이드 4: AI 에이전트란?
# ============================================================
slide = add_slide(prs)
add_title_text(slide, "AI 에이전트란?")
add_body_text(slide, """AI 에이전트 = LLM + 도구(Tool) + 판단력

일반 챗봇과의 차이:
  챗봇: 사용자 질문 -> 텍스트 답변 (끝)
  에이전트: 사용자 질문 -> 판단 -> 도구 사용 -> 결과 확인 -> 최종 답변

비유로 이해하기:
  챗봇 = "전화 상담원" (말만 할 수 있음)
  에이전트 = "현장 직원" (직접 시스템을 조작하고 결과를 확인함)

에이전트의 핵심 능력:
  1. 어떤 도구를 사용할지 스스로 판단 (LLM의 추론 능력)
  2. 도구를 실행하고 결과를 확인
  3. 필요하면 추가 도구를 사용하거나 최종 답변 생성

예시: "3+4를 구하고 결과에 7을 곱해줘"
  -> LLM이 add(3,4) 호출 판단
  -> 결과 7 확인
  -> multiply(7,7) 호출
  -> 49 답변""", font_size=13)

# ============================================================
# 슬라이드 5: LangGraph 3대 핵심 요소
# ============================================================
slide = add_slide(prs)
add_title_text(slide, "LangGraph 3대 핵심 요소")
add_table_slide(slide,
    headers=["구성 요소", "설명", "비유"],
    rows=[
        ["State (상태)", "에이전트 실행 중 모든 노드가\n공유하는 데이터 구조", "공유 메모장\n(모든 직원이 볼 수 있는 칠판)"],
        ["Node (노드)", "실제 작업을 수행하는 함수\n(LLM 호출, 도구 실행 등)", "작업자\n(각자 맡은 업무를 수행)"],
        ["Edge (엣지)", "노드 간의 연결선\n실행 순서와 조건을 정의", "연결 통로\n(다음 작업자에게 전달)"],
    ],
    top=Inches(1.8)
)
add_body_text(slide, """정리: State에 데이터를 저장하고, Node가 작업을 수행하며, Edge가 순서를 결정합니다.""",
    top=Inches(4.5), font_size=14, color=DARK_GRAY)

# ============================================================
# 슬라이드 6: 그래프 구조 이해
# ============================================================
slide = add_slide(prs)
add_title_text(slide, "그래프 구조 이해")
add_body_text(slide, """LangGraph는 방향 그래프(Directed Graph)를 사용합니다.

  방향 그래프: 화살표가 있는 그래프 (A -> B는 되지만, B -> A는 별도 설정 필요)
  모든 그래프에는 시작점(START)과 끝점(END)이 있습니다.""", font_size=14)

add_flow_diagram(slide, """START  -->  [Node A]  -->  [Node B]  -->  END

예시: START --> [LLM 호출] --> [도구 실행] --> [LLM 호출] --> END""",
    top=Inches(3.5), height=Inches(1.5))

add_body_text(slide, """두 가지 종류의 엣지:
  1. 고정 엣지: A 다음에 항상 B 실행 (add_edge)
  2. 조건부 엣지: 조건에 따라 B 또는 C로 분기 (add_conditional_edges)""",
    top=Inches(5.5), font_size=13)

# ============================================================
# 슬라이드 7: 조건부 분기란?
# ============================================================
slide = add_slide(prs)
add_title_text(slide, "조건부 분기 (Conditional Edge)")
add_body_text(slide, """프로그래밍의 if-else와 같은 개념입니다.

"LLM 응답에 도구 호출이 포함되어 있으면 도구를 실행하고,
 아니면 대화를 종료한다"

이런 판단을 그래프에서 표현하는 방법이 조건부 엣지입니다.""", font_size=14)

add_flow_diagram(slide, """                    [llm_call 노드]
                         |
                   (도구 호출 있나?)
                    /          \\
                  Yes           No
                  /               \\
          [tool_node]            [END]
              |
          [llm_call] (다시)""",
    top=Inches(3.5), height=Inches(2.5))

# ============================================================
# 슬라이드 8: Tool(도구) 개념
# ============================================================
slide = add_slide(prs)
add_title_text(slide, "Tool (도구) 개념")
add_body_text(slide, """LLM은 텍스트 생성만 가능합니다. 실제 계산/검색/API 호출은 못합니다.

도구(Tool) = LLM이 사용할 수 있는 Python 함수

작동 방식:
  1. LLM이 사용자 질문을 분석
  2. "이 질문에는 multiply 함수가 필요하다"고 판단
  3. 함수 이름과 인자를 JSON으로 반환 (tool_calls)
  4. 시스템이 실제 함수를 실행
  5. 결과를 LLM에게 다시 전달
  6. LLM이 결과를 바탕으로 최종 답변 생성""", font_size=14)

add_code_block(slide, """# @tool 데코레이터로 일반 함수를 도구로 변환
@tool
def multiply(a: int, b: int) -> int:
    \"\"\"Multiply a and b.\"\"\"   # <-- LLM이 이 설명을 읽고 사용 여부 판단
    return a * b""", top=Inches(5.0), font_size=12)

# ============================================================
# 슬라이드 9: Part 1 개요
# ============================================================
slide = add_slide(prs)
add_title_text(slide, "Part 1. 사칙연산 에이전트 구축", font_size=30, color=ACCENT)
add_body_text(slide, """목표: LangGraph의 기본 구조를 이해하고, 사칙연산 에이전트를 직접 구현

구현 단계:
  Step 1. Model 및 Tool 정의 (LLM 모델 + 사칙연산 함수)
  Step 2. State 정의 (공유 데이터 구조)
  Step 3. Node 정의 (LLM 호출 노드 + 도구 실행 노드)
  Step 4. Graph 생성 (노드와 엣지 연결)
  Step 5. 에이전트 실행""", font_size=15)

add_flow_diagram(slide, """START --> [llm_call] --> (도구 필요?) --Yes--> [tool_node] --> [llm_call] (반복)
                                        |
                                       No
                                        |
                                      [END]""",
    top=Inches(4.8), height=Inches(1.5))

# ============================================================
# 슬라이드 10: Step 1 - 모델 생성
# ============================================================
slide = add_slide(prs)
add_title_text(slide, "Step 1. 모델 생성")
add_body_text(slide, """init_chat_model() 함수로 LLM 모델을 생성합니다.

  모델 이름만 넣으면 자동으로 해당 provider를 감지합니다
  temperature=0: 일관된 응답 (창의성 없음, 수학 계산에 적합)
  temperature=1: 창의적 응답 (글쓰기 등에 적합)""", font_size=14)

add_code_block(slide, """from langchain.chat_models import init_chat_model

model = init_chat_model(
    "gpt-4o-mini",    # 사용할 모델명
    temperature=0     # 응답의 일관성을 최대화
)""", top=Inches(3.5), font_size=13)

# ============================================================
# 슬라이드 11: Step 1 - 도구 정의
# ============================================================
slide = add_slide(prs)
add_title_text(slide, "Step 1. 도구(Tool) 정의")
add_body_text(slide, """@tool 데코레이터로 일반 Python 함수를 LangChain 도구로 변환합니다.
LLM은 함수의 이름과 docstring을 읽고 언제 사용할지 스스로 판단합니다.""", font_size=14)

add_code_block(slide, """from langchain.tools import tool

@tool
def multiply(a: int, b: int) -> int:
    \"\"\"Multiply a and b.
    Args:
        a: First int
        b: Second int
    \"\"\"
    return a * b

@tool
def add(a: int, b: int) -> int:
    \"\"\"Adds a and b.\"\"\"
    return a + b

@tool
def divide(a: int, b: int) -> float:
    \"\"\"Divide a and b.\"\"\"
    return a / b

# 도구 목록 생성
tools = [multiply, add, divide]""", top=Inches(2.6), font_size=11)

# ============================================================
# 슬라이드 12: 모델+도구 바인딩
# ============================================================
slide = add_slide(prs)
add_title_text(slide, "Step 1. 모델과 도구 바인딩")
add_body_text(slide, """bind_tools()로 모델에게 "이런 도구들을 사용할 수 있다"고 알려줍니다.

  바인딩 후 모델을 호출하면, 직접 답변하거나 도구 호출을 요청합니다
  실제 도구 실행은 모델이 아니라 별도의 tool_node에서 처리합니다""", font_size=14)

add_code_block(slide, """# 모델에 도구 목록 연결
model_with_tools = model.bind_tools(tools)

# 이제 model_with_tools를 호출하면:
# - 직접 답변 가능한 질문 -> 텍스트 응답
# - 도구가 필요한 질문 -> tool_calls 포함 응답""", top=Inches(3.3), font_size=13)

add_body_text(slide, """중요 포인트:
  bind_tools()는 모델에게 도구 "정보"만 전달합니다.
  실제 함수 실행은 별도의 tool_node가 담당합니다!""",
    top=Inches(5.3), font_size=13, color=DARK_GRAY)

# ============================================================
# 슬라이드 13: Step 2 - State 정의
# ============================================================
slide = add_slide(prs)
add_title_text(slide, "Step 2. State 정의")
add_body_text(slide, """State = 에이전트 실행 중 모든 노드가 공유하는 데이터 구조

  TypedDict로 타입을 명시하여 안정성 확보
  Annotated + operator.add: 새 메시지가 오면 기존 목록에 "추가"(append)
  (덮어쓰기가 아님!)""", font_size=14)

add_code_block(slide, """from langchain.messages import AnyMessage
from typing_extensions import TypedDict, Annotated
import operator

class MessagesState(TypedDict):
    # 새 메시지는 기존 목록에 추가됨 (operator.add)
    messages: Annotated[list[AnyMessage], operator.add]

    # LLM 호출 횟수 카운터
    llm_calls: int""", top=Inches(3.3), font_size=13)

add_body_text(slide, """operator.add의 동작:
  기존: [msg1, msg2]  +  새로운: [msg3]  =  결과: [msg1, msg2, msg3]""",
    top=Inches(5.5), font_size=13, color=DARK_GRAY)

# ============================================================
# 슬라이드 14: Step 3 - llm_call 노드
# ============================================================
slide = add_slide(prs)
add_title_text(slide, "Step 3. Node 정의 (1) - llm_call")
add_body_text(slide, """llm_call 노드: LLM을 호출하여 응답을 생성하거나 도구 사용을 요청합니다.""", font_size=14)

add_code_block(slide, """from langchain.messages import SystemMessage

def llm_call(state):
    \"\"\"LLM이 답변하거나 도구 사용을 요청하는 Node\"\"\"

    # SystemMessage로 LLM에게 역할 부여
    response = model_with_tools.invoke(
        [
            SystemMessage(
                content="당신은 사칙연산을 하는 유능한 Agent입니다."
            )
        ] + state["messages"]  # 시스템 메시지 + 기존 대화 기록
    )

    return {
        "messages": [response],                      # 새 응답 추가
        "llm_calls": state.get('llm_calls', 0) + 1  # 호출 횟수 증가
    }""", top=Inches(2.2), font_size=11)

add_body_text(slide, """핵심: 노드 함수는 state를 받아서 변경된 부분만 딕셔너리로 반환합니다.
LangGraph가 자동으로 State에 병합합니다.""",
    top=Inches(6.0), font_size=13, color=DARK_GRAY)

# ============================================================
# 슬라이드 15: Step 3 - tool_node
# ============================================================
slide = add_slide(prs)
add_title_text(slide, "Step 3. Node 정의 (2) - tool_node")
add_body_text(slide, """tool_node: LLM이 요청한 도구를 실제로 실행하는 노드입니다.""", font_size=14)

add_code_block(slide, """from langchain.messages import ToolMessage

# 도구 이름 -> 함수 매핑
tools_by_name = {tool.name: tool for tool in tools}

def tool_node(state):
    \"\"\"LLM이 요청한 도구를 실행하는 Node\"\"\"
    result = []

    for tool_call in state["messages"][-1].tool_calls:
        # 1. 도구 이름으로 함수 찾기
        tool = tools_by_name[tool_call["name"]]

        # 2. 함수 실행
        tool_result = tool.invoke(tool_call["args"])

        # 3. 결과를 ToolMessage로 포장
        result.append(
            ToolMessage(
                content=tool_result,
                tool_call_id=tool_call["id"]  # 요청-응답 매칭 ID
            )
        )

    return {"messages": result}""", top=Inches(2.0), font_size=11)

# ============================================================
# 슬라이드 16: Step 4 - 그래프 생성
# ============================================================
slide = add_slide(prs)
add_title_text(slide, "Step 4. 그래프 생성")
add_body_text(slide, """StateGraph로 그래프를 만들고, 노드와 엣지를 추가합니다.""", font_size=14)

add_code_block(slide, """from langgraph.graph import StateGraph, START

# 1. 그래프 빌더 생성 (State 스키마 지정)
agent_builder = StateGraph(MessagesState)

# 2. 노드 등록
agent_builder.add_node("llm_call", llm_call)    # LLM 호출 노드
agent_builder.add_node("tool_node", tool_node)   # 도구 실행 노드

# 3. 고정 엣지 연결
agent_builder.add_edge(START, "llm_call")        # 시작 -> LLM 호출
agent_builder.add_edge("tool_node", "llm_call")  # 도구 실행 후 -> LLM 다시 호출""", top=Inches(2.2), font_size=12)

add_body_text(slide, """add_node(이름, 함수): 그래프에 노드(작업자)를 추가
add_edge(출발, 도착): 항상 실행되는 연결선을 추가
START: 그래프 실행의 시작점 (특수 상수)""",
    top=Inches(5.2), font_size=13, color=DARK_GRAY)

# ============================================================
# 슬라이드 17: Step 4 - 조건부 엣지
# ============================================================
slide = add_slide(prs)
add_title_text(slide, "Step 4. 조건부 엣지 (Conditional Edge)")
add_body_text(slide, """LLM 응답에 따라 다음 노드를 동적으로 결정합니다.""", font_size=14)

add_code_block(slide, """from langgraph.graph import END

def should_continue(state: MessagesState):
    \"\"\"다음 단계를 결정하는 라우팅 함수\"\"\"
    last_message = state['messages'][-1]

    if last_message.tool_calls:
        return "tool_node"   # 도구 사용 요청 -> tool_node 실행
    else:
        return END           # 최종 답변 완성 -> 종료

# 조건부 엣지 연결
agent_builder.add_conditional_edges(
    'llm_call',         # 출발 노드
    should_continue,    # 다음 노드를 결정하는 함수
    [END, 'tool_node']  # 가능한 목적지 목록
)""", top=Inches(2.2), font_size=12)

add_body_text(slide, """핵심: should_continue 함수가 반환하는 문자열이 다음 실행할 노드의 이름입니다.""",
    top=Inches(6.0), font_size=13, color=DARK_GRAY)

# ============================================================
# 슬라이드 18: 실행 흐름도
# ============================================================
slide = add_slide(prs)
add_title_text(slide, "Step 4. 전체 실행 흐름")
add_body_text(slide, """"3과 4를 더하고 결과에 7을 곱해줘" 실행 흐름:""", font_size=14)

add_flow_diagram(slide, """[1] START --> llm_call
    입력: "3과 4를 더하고 결과에 7을 곱해줘"
    LLM 판단: add(3,4) 호출 필요 -> tool_calls 반환

[2] llm_call --> tool_node (should_continue: tool_calls 있음)
    tool_node: add(3,4) 실행 -> 결과: 7

[3] tool_node --> llm_call (고정 엣지)
    LLM 판단: multiply(7,7) 호출 필요 -> tool_calls 반환

[4] llm_call --> tool_node
    tool_node: multiply(7,7) 실행 -> 결과: 49

[5] tool_node --> llm_call
    LLM: "결과는 49입니다" (tool_calls 없음)

[6] llm_call --> END (should_continue: tool_calls 없음)""",
    top=Inches(2.2), height=Inches(4.5))

# ============================================================
# 슬라이드 19: Step 5 - 에이전트 실행
# ============================================================
slide = add_slide(prs)
add_title_text(slide, "Step 5. 에이전트 실행")

add_code_block(slide, """# 그래프 컴파일 (에이전트 완성)
agent = agent_builder.compile()

# 에이전트 실행
from langchain.messages import HumanMessage

messages = [HumanMessage(content="3과 4를 더해줘")]
response = agent.invoke({"messages": messages})

# 최종 답변 확인
print(response["messages"][-1].content)
# 출력: "3과 4를 더한 결과는 7입니다." """, top=Inches(1.7), font_size=13)

add_body_text(slide, """invoke() 메서드:
  에이전트를 동기적으로 실행하고 최종 결과를 반환
  입력 형식: {"messages": [메시지 목록]}
  반환값: 전체 State (messages, llm_calls 등)

compile() 메서드:
  그래프 정의를 실행 가능한 에이전트로 변환
  노드/엣지 유효성 검사 수행""", top=Inches(4.8), font_size=13)

# ============================================================
# 슬라이드 20: Part 2 개요
# ============================================================
slide = add_slide(prs)
add_title_text(slide, "Part 2. 이메일 처리 에이전트 (규칙 기반)", font_size=26, color=ACCENT)
add_body_text(slide, """목표: Router 패턴을 구현하여 조건에 따라 다른 처리 경로로 분기

시나리오: 고객 이메일을 자동으로 분류하고 처리
  환불/긴급 키워드 포함 -> 상담원에게 이관
  일반 문의 -> 자동 답변 생성

핵심 패턴: Router (라우터)
  "입력을 분석해서 적절한 처리 경로로 보내는 것"
  (도로의 분기점처럼 어디로 갈지 결정)""", font_size=15)

add_flow_diagram(slide, """START --> [이메일 읽기] --> [의도 분류]
                                          |
                              (complaint? / inquiry?)
                               /                  \\
                    [상담원 이관]            [매뉴얼 검색]
                        |                       |
                      [END]               [답변 생성]
                                              |
                                            [END]""",
    top=Inches(4.5), height=Inches(2.3))

# ============================================================
# 슬라이드 21: Part 2 - State 정의
# ============================================================
slide = add_slide(prs)
add_title_text(slide, "Part 2. State 정의 - 이메일 처리용")
add_body_text(slide, """이메일 처리에 특화된 State 구조를 정의합니다.
Part 1과 달리 messages가 아닌 각 단계의 결과를 개별 필드로 관리합니다.""", font_size=14)

add_code_block(slide, """from typing_extensions import TypedDict

class AgentState(TypedDict):
    email_content: str  # 처리할 이메일 원본 내용
    category: str       # 분류 결과 ('complaint' 또는 'inquiry')
    next_step: str      # 다음 처리 단계 (라우팅에 사용)
    response: str       # 최종 AI 생성 답변""", top=Inches(2.8), font_size=13)

add_table_slide(slide,
    headers=["필드", "타입", "역할"],
    rows=[
        ["email_content", "str", "처리할 이메일 원본"],
        ["category", "str", "분류 결과 (complaint/inquiry)"],
        ["next_step", "str", "라우팅 방향 결정"],
        ["response", "str", "AI가 생성한 답변"],
    ],
    top=Inches(4.8)
)

# ============================================================
# 슬라이드 22: Part 2 - 노드 정의
# ============================================================
slide = add_slide(prs)
add_title_text(slide, "Part 2. Node 정의 - 이메일 읽기 & 분류")

add_code_block(slide, """# [Node 1] 이메일 읽기 (패스스루)
def read_email(state: AgentState):
    return {'email_content': state['email_content']}

# [Node 2] 의도 분류 (규칙 기반 - 키워드 매칭)
def classify_intent(state: AgentState):
    email = state['email_content']

    # 키워드 기반 분류
    if '환불' in email or '긴급' in email or '빨리' in email:
        category = 'complaint'           # 불만/긴급
        next_step = 'escalate_to_human'  # 상담원 이관
    else:
        category = 'inquiry'             # 일반 문의
        next_step = 'search_manual'      # 매뉴얼 검색

    return {'category': category, 'next_step': next_step}""", top=Inches(1.5), font_size=11)

add_body_text(slide, """포인트:
  read_email: 현재는 단순 전달이지만, 실무에서는 이메일 파싱/전처리 담당
  classify_intent: 키워드 매칭으로 분류 (빠르지만 한계 있음)
    - '환불'은 감지하지만 '환급'은 감지 못함 -> Part 3에서 개선""",
    top=Inches(5.3), font_size=13)

# ============================================================
# 슬라이드 23: Part 2 - 나머지 노드
# ============================================================
slide = add_slide(prs)
add_title_text(slide, "Part 2. Node 정의 - 검색, 이관, 답변 생성")

add_code_block(slide, """# [Node 3-A] 매뉴얼 검색 (자동 처리 트랙)
def search_manual(state: AgentState):
    print('매뉴얼을 검색합니다')
    return  # 실제로는 RAG 시스템과 연동

# [Node 3-B] 상담원 이관 (긴급 이슈 트랙)
def escalate_to_human(state: AgentState):
    print('상담원 이관합니다')
    return  # 실제로는 CRM 티켓 생성, 슬랙 알림 등

# [Node 4] 답변 생성
def write_reply(state: AgentState):
    email = state['email_content']
    response = model.invoke(email)  # LLM으로 답변 생성
    return {'response': response}""", top=Inches(1.5), font_size=12)

add_body_text(slide, """실무 확장 포인트:
  search_manual -> 벡터 DB 검색 (RAG) 시스템 연동
  escalate_to_human -> CRM API 호출, Slack 알림, 이메일 전송
  write_reply -> 검색 결과를 포함한 프롬프트로 더 정확한 답변 생성""",
    top=Inches(5.3), font_size=13, color=DARK_GRAY)

# ============================================================
# 슬라이드 24: Part 2 - 그래프 구성
# ============================================================
slide = add_slide(prs)
add_title_text(slide, "Part 2. 그래프 구성")

add_code_block(slide, """from langgraph.graph import StateGraph, START, END

agent_builder = StateGraph(AgentState)

# 노드 등록
agent_builder.add_node('read_email', read_email)
agent_builder.add_node('classify_intent', classify_intent)
agent_builder.add_node('search_manual', search_manual)
agent_builder.add_node('escalate_to_human', escalate_to_human)
agent_builder.add_node('write_reply', write_reply)

# 고정 엣지
agent_builder.add_edge(START, 'read_email')
agent_builder.add_edge('read_email', 'classify_intent')

# 라우팅 함수
def route_email(state: AgentState):
    if state['next_step'] == 'escalate_to_human':
        return 'escalate_to_human'
    else:
        return 'search_manual'

# 조건부 엣지 (분류 결과에 따라 분기)
agent_builder.add_conditional_edges(
    'classify_intent', route_email,
    ['escalate_to_human', 'search_manual']
)

# 나머지 고정 엣지
agent_builder.add_edge('search_manual', 'write_reply')
agent_builder.add_edge('write_reply', END)
agent_builder.add_edge('escalate_to_human', END)

agent = agent_builder.compile()""", top=Inches(1.5), font_size=10)

# ============================================================
# 슬라이드 25: Part 2 - 실행 테스트
# ============================================================
slide = add_slide(prs)
add_title_text(slide, "Part 2. 실행 테스트")

add_code_block(slide, """# 테스트 1: 일반 문의
response = agent.invoke({
    "email_content": "비밀번호 변경은 어떻게 하나요?"
})
# 결과: category='inquiry', 자동 답변 생성됨

# 테스트 2: 환불 요청
response = agent.invoke({
    "email_content": "환불해주세요. 제품이 불량입니다."
})
# 결과: category='complaint', 상담원 이관 처리

# 테스트 3: '환급' 키워드
response = agent.invoke({
    "email_content": "환급 요청합니다."
})
# 결과: category='inquiry' (!)  <- 키워드 미매칭!""", top=Inches(1.5), font_size=12)

add_body_text(slide, """문제 발견:
  '환불'은 감지하지만, '환급'은 키워드 목록에 없어서 일반 문의로 분류됨!
  -> 규칙 기반 분류의 근본적 한계""",
    top=Inches(5.5), font_size=14, color=RGBColor(180, 0, 0))

# ============================================================
# 슬라이드 26: 규칙 기반의 한계
# ============================================================
slide = add_slide(prs)
add_title_text(slide, "규칙 기반 분류의 한계")
add_body_text(slide, """규칙 기반 (키워드 매칭) 분류의 문제점:

1. 키워드를 계속 추가/관리해야 함
   - '환불', '환급', '반품', '취소', '돌려주세요'... 끝이 없음

2. 문맥을 이해하지 못함
   - "환불 정책이 궁금합니다" -> 불만이 아닌 단순 문의인데 '환불' 키워드로 이관됨

3. 새로운 표현에 대응 불가
   - "이거 돈 다시 줄 수 있나요?" -> 키워드 없음, 감지 실패

해결책: LLM 기반 분류!
  - LLM이 문맥과 의도를 직접 파악
  - 키워드 목록 없이도 다양한 표현을 정확하게 분류
  - Part 3에서 구현합니다""", font_size=14)

# ============================================================
# 슬라이드 27: Part 3 개요
# ============================================================
slide = add_slide(prs)
add_title_text(slide, "Part 3. LLM 기반 이메일 분류 에이전트", font_size=26, color=ACCENT)
add_body_text(slide, """목표: 규칙 기반 분류의 한계를 LLM으로 극복

핵심 변경점:
  classify_intent (키워드 매칭) -> classify_node (LLM 판단)

LLM 기반 분류의 장점:
  키워드 목록 관리 불필요
  문맥과 의도를 이해하여 더 정확한 분류
  '환급', '돈 돌려줘' 등 다양한 표현 처리 가능

추가 개선:
  messages 기반 State로 대화 기록 관리
  도구(Tool) 활용 상담 AI (consultant_node)""", font_size=15)

# ============================================================
# 슬라이드 28: 규칙 vs LLM 비교
# ============================================================
slide = add_slide(prs)
add_title_text(slide, "규칙 기반 vs LLM 기반 분류 비교")

add_table_slide(slide,
    headers=["항목", "규칙 기반 (Part 2)", "LLM 기반 (Part 3)"],
    rows=[
        ["분류 방법", "키워드 매칭 (if-else)", "LLM이 문맥 분석"],
        ["'환불' 처리", "O (키워드 매칭)", "O (의미 이해)"],
        ["'환급' 처리", "X (키워드 없음)", "O (유사 의미 이해)"],
        ["문맥 이해", "X", "O"],
        ["유지보수", "키워드 계속 추가", "프롬프트만 관리"],
        ["비용", "무료 (코드 실행)", "API 호출 비용 발생"],
        ["속도", "매우 빠름", "LLM 호출 시간 필요"],
    ],
    top=Inches(1.8)
)

add_body_text(slide, """선택 기준: 정확도가 중요하면 LLM 기반, 속도/비용이 중요하면 규칙 기반
실무에서는 두 방식을 조합하여 사용하기도 합니다.""",
    top=Inches(5.5), font_size=13, color=DARK_GRAY)

# ============================================================
# 슬라이드 29: classify_node 코드
# ============================================================
slide = add_slide(prs)
add_title_text(slide, "Part 3. classify_node - LLM 기반 분류")
add_body_text(slide, """System Prompt로 LLM에게 역할과 출력 형식을 지정합니다.""", font_size=14)

add_code_block(slide, """def classify_node(state: AgentState):
    \"\"\"LLM을 사용하여 이메일 의도를 분류\"\"\"
    last_message = state["messages"][-1]

    # 분류용 System Prompt
    prompt = \"\"\"
    당신은 고객 센터 관리자입니다.
    고객의 이메일을 분석해서 다음 단계를 결정하세요.

    1. 단순 문의나 정보 요청 -> 'consultant' 반환
    2. 환불 요청, 불만 제기, 화난 고객 -> 'escalate' 반환

    답변은 오직 단어 하나만 하세요.
    \"\"\"

    response = model.invoke(
        [SystemMessage(content=prompt)] + [last_message]
    )
    decision = str(response.content).strip().lower()

    if "escalate" in decision:
        return {"next_step": "escalate"}
    else:
        return {"next_step": "consultant"}""", top=Inches(1.9), font_size=11)

# ============================================================
# 슬라이드 30: consultant_node + tool_node
# ============================================================
slide = add_slide(prs)
add_title_text(slide, "Part 3. consultant_node & tool_node")
add_body_text(slide, """상담 AI가 도구(매뉴얼 검색)를 활용하여 고객에게 답변합니다.""", font_size=14)

add_code_block(slide, """# 매뉴얼 검색 도구
@tool
def search_manual(query: str) -> str:
    \"\"\"고객 질문에 답변하기 위한 매뉴얼 검색 도구\"\"\"
    if '비밀번호' in query:
        return '비밀번호 변경은 마이페이지 - 보안 설정에 있음'
    elif '배송' in query:
        return '00택배에서 3일 내 배송 예정임'
    else:
        return '해당 내용 관련 매뉴얼은 찾을 수 없습니다.'

model_with_tools = model.bind_tools([search_manual])

# 상담 AI 노드
def consultant_node(state: AgentState):
    response = model_with_tools.invoke(state['messages'])
    return {'messages': [response]}

# 이관 노드
def escalate_node(state: AgentState):
    return {'messages': [
        AIMessage(content='해당 메일은 전문 상담원에게 이관되었습니다.')
    ]}""", top=Inches(2.0), font_size=11)

# ============================================================
# 슬라이드 31: Part 3 - 그래프 구성
# ============================================================
slide = add_slide(prs)
add_title_text(slide, "Part 3. 그래프 구성 - 이중 조건부 분기")

add_code_block(slide, """agent_builder = StateGraph(AgentState)

# 노드 등록
agent_builder.add_node('classify_node', classify_node)
agent_builder.add_node('consultant_node', consultant_node)
agent_builder.add_node('escalate_node', escalate_node)
agent_builder.add_node('tool_node', tool_node)

# 시작 엣지
agent_builder.add_edge(START, 'classify_node')

# 조건부 엣지 1: 분류 결과에 따라 분기
agent_builder.add_conditional_edges(
    'classify_node', route_after_classify,
    {'escalate': 'escalate_node', 'consultant': 'consultant_node'}
)

# 조건부 엣지 2: 도구 사용 여부 판단
agent_builder.add_conditional_edges(
    'consultant_node', should_continue,
    ['tool_node', END]
)

# 고정 엣지
agent_builder.add_edge('tool_node', 'consultant_node')
agent_builder.add_edge('escalate_node', END)""", top=Inches(1.5), font_size=11)

add_flow_diagram(slide, """START --> [classify_node] --escalate--> [escalate_node] --> END
                       |
                   consultant
                       |
               [consultant_node] --tool_calls--> [tool_node] --> [consultant_node]
                       |
                    (no tools)
                       |
                     [END]""",
    top=Inches(5.3), height=Inches(1.8))

# ============================================================
# 슬라이드 32: Part 3 - 실행 결과 비교
# ============================================================
slide = add_slide(prs)
add_title_text(slide, "Part 3. 실행 결과 - LLM 기반 분류의 위력")

add_code_block(slide, """# 테스트 1: 일반 문의 ("비밀번호 변경 방법")
response = agent.invoke({
    "messages": [HumanMessage(content="비밀번호 변경은 어떻게 하나요?")]
})
# -> classify_node: 'consultant' 판단
# -> consultant_node: search_manual 도구 사용
# -> 최종 답변: "비밀번호 변경은 마이페이지 - 보안 설정에서..."

# 테스트 2: 환불 요청 ("환급" - Part 2에서 실패했던 키워드!)
response = agent.invoke({
    "messages": [HumanMessage(content="환급 요청합니다.")]
})
# -> classify_node: 'escalate' 판단 (LLM이 의미를 이해!)
# -> escalate_node: "전문 상담원에게 이관되었습니다"
# -> 성공! 규칙 기반에서 실패했던 케이스 해결!""", top=Inches(1.5), font_size=12)

add_body_text(slide, """핵심 성과:
  '환급'이라는 키워드를 직접 등록하지 않았는데도
  LLM이 "환불과 유사한 의미"라고 판단하여 정확하게 이관 처리!

  이것이 LLM 기반 분류의 핵심 장점입니다.""",
    top=Inches(5.5), font_size=14, color=RGBColor(0, 120, 0))

# ============================================================
# 슬라이드 33: State 설계 원칙
# ============================================================
slide = add_slide(prs)
add_title_text(slide, "State 설계 원칙")
add_body_text(slide, """좋은 State 설계를 위한 핵심 원칙:

1. 필요한 것만 포함
   - 모든 노드가 공유하는 정보만 State에 넣습니다
   - 한 노드에서만 쓰는 임시 데이터는 지역 변수로 처리

2. 명확한 타입 정의
   - TypedDict로 각 필드의 타입을 명시
   - 코드 안정성과 가독성 향상

3. 누적 vs 덮어쓰기 결정
   - messages처럼 쌓여야 하는 데이터: Annotated + operator.add
   - 최신값만 필요한 데이터: 일반 타입 (덮어쓰기)

4. 라우팅 정보 포함
   - next_step 같은 필드로 조건부 분기를 제어
   - 노드가 State에 결과를 쓰고, 라우팅 함수가 읽어서 분기

Part 1 State (메시지 중심):
  messages + llm_calls

Part 2 State (이메일 처리 특화):
  email_content + category + next_step + response

Part 3 State (대화 + 라우팅):
  messages + next_step""", font_size=13)

# ============================================================
# 슬라이드 34: 핵심 요약
# ============================================================
slide = add_slide(prs)
add_title_text(slide, "핵심 요약", font_size=30)

add_body_text(slide, """LangGraph 3대 요소:
  State = 공유 데이터 | Node = 작업 함수 | Edge = 실행 순서

에이전트 구축 5단계:
  1. Model & Tool 정의
  2. State 정의
  3. Node 정의
  4. Graph 생성 (노드 + 엣지 연결)
  5. 컴파일 & 실행

두 가지 엣지:
  고정 엣지 (add_edge): 항상 같은 경로
  조건부 엣지 (add_conditional_edges): 조건에 따라 분기

분류 방식 비교:
  규칙 기반: 빠르고 저렴하지만, 키워드 관리 필요 + 문맥 이해 불가
  LLM 기반: 느리고 비용 발생하지만, 의미 이해 + 유연한 분류

핵심 패턴:
  Tool Use 패턴: LLM -> 도구 호출 -> 결과 반환 -> 최종 답변
  Router 패턴: 조건에 따라 다른 처리 경로로 분기""", font_size=13)

# ============================================================
# 슬라이드 35: Q&A
# ============================================================
slide = add_slide(prs)
add_title_text(slide, "Q & A", top=Inches(2.8), font_size=40, color=BLACK)
add_subtitle_text(slide, "질문이 있으시면 말씀해 주세요!", top=Inches(3.6), font_size=20, color=DARK_GRAY)

# ============================================================
# 저장
# ============================================================
output_path = r"C:\DL_Practice\2강_랭그래프_기초_강의슬라이드.pptx"
prs.save(output_path)
print(f"PPT 파일 생성 완료: {output_path}")
print(f"총 슬라이드 수: {len(prs.slides)}")
